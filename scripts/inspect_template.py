"""Inspect 5주차 PPT template — list slide layouts and master info."""
from pptx import Presentation
from pptx.util import Emu

TEMPLATE = r"C:\Users\kk409\Documents\Konkuk\3-1\Unreal_Study\언리얼_스터디_5주차_자료_수정본.pptx"

prs = Presentation(TEMPLATE)
print(f"Slide size: {Emu(prs.slide_width).inches:.2f}in x {Emu(prs.slide_height).inches:.2f}in")
print(f"Slide masters: {len(prs.slide_masters)}")
for mi, master in enumerate(prs.slide_masters):
    print(f"\n--- Master {mi} ---")
    print(f"  Layouts: {len(master.slide_layouts)}")
    for li, layout in enumerate(master.slide_layouts):
        ph_summary = []
        for ph in layout.placeholders:
            ph_summary.append(f"idx={ph.placeholder_format.idx} type={ph.placeholder_format.type} name={ph.name!r}")
        print(f"  [{li}] name={layout.name!r}")
        for s in ph_summary:
            print(f"        {s}")

print("\n--- Existing slides (first 3) ---")
for i, slide in enumerate(prs.slides):
    if i >= 3:
        break
    layout = slide.slide_layout
    print(f"Slide {i+1}: layout={layout.name!r}")
    for ph in slide.placeholders:
        text = ph.text_frame.text[:50] if ph.has_text_frame else ""
        print(f"  ph idx={ph.placeholder_format.idx} text={text!r}")

# Sample run text & font info from slide 2 (a content slide)
print("\n--- Font samples from slide 2 text runs ---")
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip():
                f = run.font
                print(f"  text={run.text[:30]!r} name={f.name} size={f.size} bold={f.bold} color_type={f.color.type if f.color and f.color.type else None}")
                break
        break
    break
