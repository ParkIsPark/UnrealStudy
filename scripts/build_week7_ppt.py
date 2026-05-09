"""Build 언리얼_스터디_7주차_자료.pptx from 5주차 template (16:9, 13.33x7.5in).

Strategy:
- Open 5주차 .pptx as template (theme/master inherited)
- Strip all existing slides
- Add 50 new slides per the slide-by-slide spec
- All slides use the single 'DEFAULT' layout, with text boxes positioned manually

Tone & layout: lifted verbatim from build_week6_ppt.py (same helper library,
same color palette, same 1.8 min/slide cadence).
"""
import io
import sys
from copy import deepcopy

# Force UTF-8 stdout for Windows cp949 console
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

TEMPLATE = r"C:\Users\kk409\Documents\Konkuk\3-1\Unreal_Study\언리얼_스터디_5주차_자료_수정본.pptx"
OUTPUT = r"C:\Users\kk409\Documents\Konkuk\3-1\Unreal_Study\언리얼_스터디_7주차_자료.pptx"

# Fonts
FONT_KR = "맑은 고딕"
FONT_CODE = "Consolas"

# Colors
COLOR_TITLE = RGBColor(0x1F, 0x1F, 0x1F)        # near-black for titles
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)         # body text
COLOR_ACCENT = RGBColor(0x2B, 0x57, 0x9A)       # blue accent
COLOR_STAR = RGBColor(0xC0, 0x39, 0x2B)         # red-ish for ★ insight
COLOR_CHECK = RGBColor(0x2E, 0x7D, 0x32)        # green for ✓
COLOR_CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)      # light gray for code blocks
COLOR_TABLE_HEADER_BG = RGBColor(0x2B, 0x57, 0x9A)
COLOR_TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TABLE_ROW_ALT = RGBColor(0xF0, 0xF4, 0xF8)
COLOR_DIVIDER_BG = RGBColor(0x1F, 0x1F, 0x1F)
COLOR_DIVIDER_FG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_QUOTE_BG = RGBColor(0xFF, 0xF8, 0xE1)

# Slide region constants (16:9, 13.33x7.5in)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)


def remove_all_slides(prs):
    """Remove every slide so we start clean."""
    sldIdLst = prs.slides._sldIdLst
    rId_list = [sld.attrib[qn('r:id')] for sld in list(sldIdLst)]
    for rId in rId_list:
        prs.part.drop_rel(rId)
    for sld in list(sldIdLst):
        sldIdLst.remove(sld)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[0])


def set_para_font(para, name=FONT_KR, size=18, bold=False, italic=False, color=COLOR_BODY):
    for run in para.runs:
        run.font.name = name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        # Also set East Asian font in rPr
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
            rPr.append(ea)
        else:
            ea.set('typeface', FONT_KR)


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, italic=False,
                color=COLOR_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_KR,
                fill=None, line_color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if line_color is not None:
        tb.line.color.rgb = line_color
    else:
        tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        set_para_font(p, name=font, size=size, bold=bold, italic=italic, color=color)
    return tb


def add_title(slide, text, *, size=32, color=COLOR_TITLE):
    return add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.85),
        text, size=size, bold=True, color=color, align=PP_ALIGN.LEFT,
    )


def add_underline(slide):
    """Thin accent line under the title."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.18), Inches(12.33), Emu(20000))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()


def add_star_footer(slide, text):
    add_textbox(
        slide, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.6),
        f"★  {text}", size=16, italic=True, bold=True, color=COLOR_STAR,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        fill=RGBColor(0xFF, 0xF4, 0xE0),
    )


def add_page_number(slide, n, total=50):
    add_textbox(
        slide, Inches(12.33), Inches(7.05), Inches(0.95), Inches(0.35),
        f"{n} / {total}", size=10, color=RGBColor(0x99, 0x99, 0x99),
        align=PP_ALIGN.RIGHT,
    )


# ----------- Slide builders -----------

def slide_title_page(prs, title, subtitle, footer):
    s = blank_slide(prs)
    add_textbox(s, Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.4),
                title, size=54, bold=True, color=COLOR_TITLE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(3.9), Inches(12.33), Inches(0.9),
                subtitle, size=28, italic=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.6),
                footer, size=18, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER, font=FONT_CODE)
    add_textbox(s, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.4),
                "박준건  |  건국대 언리얼 스터디  |  2026-05-18", size=12,
                color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)
    return s


def slide_divider(prs, n, title, subtitle):
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_DIVIDER_BG
    bg.line.fill.background()
    add_textbox(s, Inches(0.5), Inches(2.6), Inches(12.33), Inches(1.6),
                title, size=60, bold=True, color=COLOR_DIVIDER_FG, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(4.4), Inches(12.33), Inches(0.8),
                subtitle, size=24, italic=True, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    return s


def slide_quote(prs, n, title, quote, subtext):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.0), Inches(2.3), Inches(11.33), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_QUOTE_BG
    box.line.color.rgb = COLOR_ACCENT
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f'"{quote}"'
    set_para_font(p, size=36, bold=True, color=COLOR_TITLE)
    add_textbox(s, Inches(1.0), Inches(4.7), Inches(11.33), Inches(2.0),
                subtext, size=18, italic=True, color=COLOR_BODY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    add_page_number(s, n)
    return s


def slide_bullets(prs, n, title, bullets, *, star=None, sub_top=1.5, body_size=20):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    body_top = Inches(sub_top)
    body_h = Inches(7.0 - sub_top - (0.7 if star else 0.2))
    tb = s.shapes.add_textbox(Inches(0.7), body_top, Inches(12.0), body_h)
    tb.fill.background()
    tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        text = b
        if text.startswith("✓ "):
            run = p.add_run()
            run.text = "✓  "
            run.font.name = FONT_KR
            run.font.size = Pt(body_size)
            run.font.bold = True
            run.font.color.rgb = COLOR_CHECK
            text = text[2:]
        elif text.startswith("★ "):
            run = p.add_run()
            run.text = "★  "
            run.font.name = FONT_KR
            run.font.size = Pt(body_size)
            run.font.bold = True
            run.font.color.rgb = COLOR_STAR
            text = text[2:]
        else:
            run = p.add_run()
            run.text = "•  "
            run.font.name = FONT_KR
            run.font.size = Pt(body_size)
            run.font.color.rgb = COLOR_ACCENT
        run = p.add_run()
        run.text = text
        run.font.name = FONT_KR
        run.font.size = Pt(body_size)
        run.font.color.rgb = COLOR_BODY
        for r in p.runs:
            rPr = r._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
                rPr.append(ea)
        p.space_after = Pt(8)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def style_table_cell(cell, text, *, size=14, bold=False, color=COLOR_BODY, fill=None,
                     align=PP_ALIGN.LEFT, font=FONT_KR):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
        rPr.append(ea)


def slide_table(prs, n, title, headers, rows, *, star=None, col_widths=None, body_size=14):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_top = Inches(1.55)
    table_h = Inches(min(5.0, 0.45 * n_rows + 0.2))
    tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(0.5), table_top, Inches(12.33), table_h)
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for ci, h in enumerate(headers):
        style_table_cell(tbl.cell(0, ci), h, size=body_size + 2, bold=True,
                         color=COLOR_TABLE_HEADER_FG, fill=COLOR_TABLE_HEADER_BG,
                         align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        fill = COLOR_TABLE_ROW_ALT if ri % 2 == 0 else None
        for ci, cell_text in enumerate(row):
            style_table_cell(tbl.cell(ri + 1, ci), str(cell_text), size=body_size, fill=fill)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_two_col(prs, n, title, left_title, left_body, right_title, right_body, *, star=None):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    for col_i, (sub, body, color) in enumerate([
        (left_title, left_body, COLOR_ACCENT),
        (right_title, right_body, COLOR_STAR),
    ]):
        left = Inches(0.5 + col_i * 6.4)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.55), Inches(6.0), Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        box.line.color.rgb = color
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)
        tf.margin_bottom = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = sub
        run.font.name = FONT_KR
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = color
        for line in body.split('\n'):
            pp = tf.add_paragraph()
            pp.alignment = PP_ALIGN.LEFT
            r = pp.add_run()
            r.text = line
            r.font.name = FONT_KR
            r.font.size = Pt(16)
            r.font.color.rgb = COLOR_BODY
            pp.space_after = Pt(6)
            rPr = r._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
                rPr.append(ea)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_code(prs, n, title, code, *, comment=None, star=None, code_size=14):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.55), Inches(12.33), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CODE_BG
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    box.line.width = Pt(0.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    lines = code.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_CODE
        run.font.size = Pt(code_size)
        if any(m in line for m in ["UCLASS", "UPROPERTY", "UFUNCTION", "USTRUCT", "GENERATED_BODY",
                                    "DECLARE_DELEGATE", "DECLARE_MULTICAST_DELEGATE",
                                    "DECLARE_DYNAMIC_MULTICAST_DELEGATE"]):
            run.font.color.rgb = COLOR_ACCENT
            run.font.bold = True
        else:
            run.font.color.rgb = COLOR_TITLE
    if comment:
        add_textbox(s, Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.6),
                    comment, size=14, italic=True, color=RGBColor(0x66, 0x66, 0x66),
                    align=PP_ALIGN.LEFT)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_experiment(prs, n, title, before, after, prediction):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    add_textbox(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.4),
                "Before", size=18, bold=True, color=COLOR_ACCENT)
    add_textbox(s, Inches(6.83), Inches(1.55), Inches(6.0), Inches(0.4),
                "After (이걸 떼어보면…)", size=18, bold=True, color=COLOR_STAR)
    for col_i, code in enumerate([before, after]):
        left = Inches(0.5 + col_i * 6.33)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), Inches(6.0), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CODE_BG
        box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        for i, line in enumerate(code.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line if line else " "
            run.font.name = FONT_CODE
            run.font.size = Pt(14)
            run.font.color.rgb = COLOR_TITLE
    pred = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(4.9), Inches(12.33), Inches(1.5))
    pred.fill.solid()
    pred.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    pred.line.color.rgb = COLOR_STAR
    pred.line.width = Pt(1.5)
    tf = pred.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "🤔  예측 — "
    r.font.name = FONT_KR
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = COLOR_STAR
    r2 = p.add_run()
    r2.text = prediction
    r2.font.name = FONT_KR
    r2.font.size = Pt(20)
    r2.font.italic = True
    r2.font.color.rgb = COLOR_TITLE
    rPr = r2._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
        rPr.append(ea)
    add_page_number(s, n)
    return s


def slide_experiment_result(prs, n, title, observation, meaning, star):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    boxes = [
        ("관찰 (실제 일어난 일)", observation, COLOR_ACCENT, RGBColor(0xE3, 0xF2, 0xFD)),
        ("의미 (왜 이렇게 되는가)", meaning, COLOR_STAR, RGBColor(0xFF, 0xEB, 0xEE)),
    ]
    for i, (label, body, line_color, fill_color) in enumerate(boxes):
        top = Inches(1.6 + i * 2.4)
        add_textbox(s, Inches(0.5), top, Inches(3.5), Inches(0.45),
                    label, size=18, bold=True, color=line_color)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(4.1), top, Inches(8.73), Inches(2.0))
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
        box.line.color.rgb = line_color
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = body
        r.font.name = FONT_KR
        r.font.size = Pt(18)
        r.font.color.rgb = COLOR_TITLE
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
            rPr.append(ea)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_flow(prs, n, title, flow_lines, star):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.5), Inches(1.7), Inches(10.33), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CODE_BG
    box.line.color.rgb = COLOR_ACCENT
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    for i, line in enumerate(flow_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.name = FONT_CODE
        r.font.size = Pt(18)
        r.font.color.rgb = COLOR_TITLE
        if line.strip() in ("↓", "▼"):
            r.font.color.rgb = COLOR_ACCENT
            r.font.bold = True
        p.space_after = Pt(4)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_callout(prs, n, title, body_lines, metaphor, star=None):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    body_text = "\n".join(body_lines)
    add_textbox(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(3.5),
                body_text, size=22, color=COLOR_BODY, align=PP_ALIGN.LEFT)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.5), Inches(5.2), Inches(10.33), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    box.line.color.rgb = COLOR_STAR
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "💡  비유 — "
    r.font.name = FONT_KR
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = COLOR_STAR
    r2 = p.add_run()
    r2.text = metaphor
    r2.font.name = FONT_KR
    r2.font.size = Pt(20)
    r2.font.italic = True
    r2.font.color.rgb = COLOR_TITLE
    rPr = r2._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
        rPr.append(ea)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_four_box(prs, n, title, boxes, star=None):
    """boxes = [(label, body_text), ...] up to 4."""
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    palette = [
        (COLOR_ACCENT, RGBColor(0xE3, 0xF2, 0xFD)),
        (COLOR_STAR, RGBColor(0xFF, 0xEB, 0xEE)),
        (COLOR_CHECK, RGBColor(0xE8, 0xF5, 0xE9)),
        (RGBColor(0x80, 0x57, 0xC2), RGBColor(0xF3, 0xE5, 0xF5)),
    ]
    box_w = Inches(6.0)
    box_h = Inches(2.3)
    positions = [(0.5, 1.6), (6.83, 1.6), (0.5, 4.1), (6.83, 4.1)]
    for i, (label, body) in enumerate(boxes[:4]):
        x, y = positions[i]
        line_color, fill_color = palette[i]
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
        box.line.color.rgb = line_color
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = label
        r.font.name = FONT_KR
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = line_color
        for line in body.split('\n'):
            pp = tf.add_paragraph()
            rr = pp.add_run()
            rr.text = line
            rr.font.name = FONT_KR
            rr.font.size = Pt(15)
            rr.font.color.rgb = COLOR_BODY
            pp.space_after = Pt(2)
            rPr = rr._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
                rPr.append(ea)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


# ===================== SLIDE CONTENT =====================

def build_slides(prs):
    n = 0

    # ===== A — 인트로 / 복습 (slides 1–6) =====

    # ----- Slide 1: Title -----
    n += 1
    slide_title_page(
        prs,
        title="언리얼 엔진 스터디 7주차",
        subtitle="변할 때만 알린다 — 델리게이트와 컴포넌트",
        footer="DECLARE_DYNAMIC_MULTICAST_DELEGATE  ·  UActorComponent  ·  UHealthComponent",
    )

    # ----- Slide 2: 복습 1 — 6주차 매크로 -----
    n += 1
    slide_table(
        prs, n,
        title="복습 1 — 6주차 매크로 한 글자가 결정한 것",
        headers=["매크로", "리플렉션 등록", "BP 노출", "효과"],
        rows=[
            ["UCLASS", "클래스", "Spawnable", "에디터가 클래스를 보고, BP에서 상속 가능"],
            ["UPROPERTY", "필드", "EditAnywhere / VisibleAnywhere", "디테일 패널에 뜨고, GC가 추적"],
            ["UFUNCTION", "함수", "BlueprintCallable / Implementable / Native", "BP에서 호출/오버라이드 가능"],
            ["USTRUCT", "구조체", "BlueprintType", "BP 변수로 쓸 수 있는 데이터 묶음"],
        ],
        col_widths=[2.5, 2.5, 3.5, 4.83],
        star="매크로 하나가 빠지면 — 디테일 패널에서 사라지고, GC도 안 잡고, BP에서 안 보인다",
    )

    # ----- Slide 3: 복습 2 — Native vs Implementable -----
    n += 1
    slide_flow(
        prs, n,
        title="복습 2 — 6주차 BlueprintNativeEvent 흐름",
        flow_lines=[
            "APickupItem::OnPickedUp(Native)  —  C++ 기본 동작",
            "↓",
            "BP_HealthPotion이 OnPickedUp 오버라이드  —  포션 효과",
            "↓",
            "OnPlayPickupEffect(Implementable)  —  디자이너가 사운드/파티클",
            "↓",
            "Destroy()",
        ],
        star="Native = 기본 + 확장 / Implementable = 빈 골격, 디자이너가 채운다",
    )

    # ----- Slide 4: 이번 주의 질문 -----
    n += 1
    slide_quote(
        prs, n,
        title="이번 주의 질문 (한 줄로)",
        quote="PlayerState가 데미지를 받았을 때, HUD는 그걸 어떻게 알까?",
        subtext="4주차에는 매 프레임 읽었다. (Percent Binding은 사실 매 프레임 호출되는 함수다.)\n"
                "이번 주는 반대로 — 변할 때만 알리는 구조로 바꾼다.\n"
                "그리고, 같은 HP 로직을 다음 주에 \"적\"도 똑같이 쓰게 만든다.",
    )

    # ----- Slide 5: 오늘의 목표 -----
    n += 1
    slide_bullets(
        prs, n,
        title="오늘의 목표",
        bullets=[
            "✓ 델리게이트 3종(Single / Multi / Dynamic Multi)을 구분하고 적절히 선택할 수 있다",
            "✓ 4주차 PlayerState의 HP 로직을 UHealthComponent로 추출한다",
            "✓ HUD가 매 프레임 폴링하던 방식에서 델리게이트 구독 방식으로 전환한다",
            "✓ 다음 주 \"적\"이 등장했을 때 같은 컴포넌트를 그대로 붙여 쓸 수 있는 구조를 만든다",
            "★ 한 줄로 — 폴링(Pull) 끝, 푸시(Push) 시작. 상속 끝, 조합 시작.",
        ],
        body_size=20,
    )

    # ----- Slide 6: Divider — 이론 시작 -----
    n += 1
    slide_divider(
        prs, n,
        title="이론 — 폴링은 왜 한계에 부딪치는가",
        subtitle="35분 (델리게이트 3종 + 컴포넌트 설계)",
    )

    # ===== B — 이론 1: 폴링의 한계 (slides 7–11) =====

    # ----- Slide 7: 4주차 HUD 코드 -----
    n += 1
    slide_code(
        prs, n,
        title="4주차 HUD를 다시 보자 — Percent Binding의 본질",
        code=(
            "// WBP_HpBar 의 ProgressBar.Percent 바인딩 (4주차)\n"
            "//\n"
            "// UFUNCTION(BlueprintCallable)\n"
            "float GetHpPercent()\n"
            "{\n"
            "    AStudyPlayerState* PS = GetOwningPlayer()->GetPlayerState<AStudyPlayerState>();\n"
            "    if (!PS) return 1.f;\n"
            "    return PS->CurrentHP / PS->MaxHP;\n"
            "}\n"
            "\n"
            "// 이 함수가 — 매 프레임 호출됩니다.\n"
            "// 안 변해도, 안 움직여도, 안 보고 있어도 — 매 프레임."
        ),
        comment="Percent Binding은 \"호출되는 함수\"다. UI는 함수의 반환값을 매 프레임 다시 그린다.",
        star="HUD가 PS의 타입을 직접 안다 — 강한 결합. 게다가 매 프레임 부른다 — 비싸고 무의미",
    )

    # ----- Slide 8: 실험 — Print String -----
    n += 1
    slide_experiment(
        prs, n,
        title="실험 — GetHpPercent에 Print String을 박아보자",
        before=(
            "// Before — 그냥 반환만 함\n"
            "float GetHpPercent()\n"
            "{\n"
            "    return PS->CurrentHP / PS->MaxHP;\n"
            "}"
        ),
        after=(
            "// After — Print String 추가\n"
            "float GetHpPercent()\n"
            "{\n"
            "    UE_LOG(LogTemp, Log, TEXT(\"GetHpPercent\"));\n"
            "    return PS->CurrentHP / PS->MaxHP;\n"
            "}\n"
            "\n"
            "// PIE 1초 동안 몇 번 찍힐까?"
        ),
        prediction="60 FPS면 1초에 60번 찍힐 것이다. HP가 1초에 한 번만 변해도, 호출은 60번.",
    )

    # ----- Slide 9: 실험 결과 -----
    n += 1
    slide_experiment_result(
        prs, n,
        title="실험 결과 — \"안 변해도 부른다\"",
        observation="PIE 1초 동안 약 60번 로그가 찍힘. HP가 변하지 않는 동안에도 계속.",
        meaning="UI 시스템은 \"값이 변했는지\" 모른다. 그래서 매 프레임 다시 묻는다 (= 폴링).",
        star="폴링은 안 변해도 부른다. 값을 바라보는 사람이 늘어날 때마다 비용이 곱해진다",
    )

    # ----- Slide 10: 듣는 쪽이 늘어나면 -----
    n += 1
    slide_bullets(
        prs, n,
        title="\"듣는 쪽이 늘어나면\" — 폴링은 폭발한다",
        bullets=[
            "미니맵에도 HP를 표시? → 미니맵도 매 프레임 GetHpPercent",
            "사운드 시스템(HP 30% 미만 경고)? → 또 매 프레임 비교",
            "컨트롤러 진동(HP 위험 시 떨림)? → 또 매 프레임 체크",
            "포스트 프로세스(피 묻은 화면 효과)? → 또 매 프레임",
            "★ 듣는 쪽 N명 × 매 프레임 호출 = 1초에 60N번. 그런데 HP는 1초에 1번 변할까 말까.",
        ],
        star="폴링은 듣는 쪽이 늘어날수록 비용이 선형으로 증가한다",
    )

    # ----- Slide 11: 폴링 vs 푸시 -----
    n += 1
    slide_two_col(
        prs, n,
        title="폴링(Pull) vs 푸시(Push) — 두 패러다임",
        left_title="폴링 (Pull)",
        left_body=(
            "• 듣는 쪽이 매 프레임 묻는다\n"
            "• \"HP 얼마야?\"  \"HP 얼마야?\"\n"
            "• 비용 = 듣는 쪽 수 × 프레임 수\n"
            "• 듣는 쪽이 PS 타입을 알아야 함\n"
            "• 안 변해도 호출됨 — 무의미한 일"
        ),
        right_title="푸시 (Push) — 델리게이트",
        right_body=(
            "• 변하는 쪽이 \"변했다!\" 알린다\n"
            "• 듣는 쪽은 한 번 구독하고 기다림\n"
            "• 비용 = 실제 변경 횟수 × 듣는 쪽 수\n"
            "• 듣는 쪽은 HP 값 외엔 몰라도 됨\n"
            "• 안 변하면 0번 호출 — 효율적"
        ),
        star="\"변할 때만 알리는\" 구조 — 이게 이번 주의 핵심 한 줄",
    )

    # ===== C — 이론 2: 델리게이트 3종 (slides 12–17) =====

    # ----- Slide 12: 델리게이트 정의 -----
    n += 1
    slide_callout(
        prs, n,
        title="델리게이트 — 한 줄 정의",
        body_lines=[
            "•  델리게이트 = \"함수 포인터\" + \"UE 리플렉션\"",
            "•  나중에 호출할 함수의 \"주소\"를 들고 있다가, 적절한 시점에 부른다",
            "•  C++의 std::function 과 비슷하지만 — UObject 라이프사이클과 BP 노출까지 챙긴다",
            "",
            "•  쓰임새: \"이벤트가 일어나면 알려줘\" 패턴 전체. 클릭, 충돌, HP 변경, 죽음, 로딩 완료, …",
        ],
        metaphor="구독(subscribe)과 같다 — 유튜브 채널을 구독해두면, 영상이 올라올 때만 알림이 온다",
    )

    # ----- Slide 13: Single-cast -----
    n += 1
    slide_table(
        prs, n,
        title="① Single-cast — 1:1 콜백",
        headers=["항목", "내용"],
        rows=[
            ["매크로", "DECLARE_DELEGATE / DECLARE_DELEGATE_OneParam ..."],
            ["바인딩", "딱 1개 함수만 (BindUObject / BindLambda)"],
            ["호출", ".Execute() 또는 .ExecuteIfBound()"],
            ["BP 노출", "❌ — C++ 전용"],
            ["대표 용도", "비동기 로딩 완료 콜백, \"끝나면 이거 한 번 불러줘\""],
            ["예시", "FStreamableDelegate OnLoaded — 에셋 로드 완료 시 1번 호출"],
        ],
        col_widths=[3.0, 9.83],
        star="\"한 번만, 한 명만\" — 가장 단순한 형태. 게임플레이에는 잘 안 쓴다",
    )

    # ----- Slide 14: Multi-cast -----
    n += 1
    slide_table(
        prs, n,
        title="② Multi-cast — 1:N, C++ 전용",
        headers=["항목", "내용"],
        rows=[
            ["매크로", "DECLARE_MULTICAST_DELEGATE / ..._OneParam ..."],
            ["바인딩", "여러 함수 가능 (AddUObject / AddLambda / AddRaw)"],
            ["호출", ".Broadcast() — 모든 구독자에게 한 번에"],
            ["BP 노출", "❌ — C++ 전용 (직렬화 불가)"],
            ["대표 용도", "엔진 내부 이벤트 (FCoreDelegates 류)"],
            ["예시", "FCoreDelegates::OnPostEngineInit — 엔진 초기화 후 호출"],
        ],
        col_widths=[3.0, 9.83],
        star="C++끼리만 통신할 때. BP가 들을 일이 없을 때만 — 우리는 거의 안 쓴다",
    )

    # ----- Slide 15: Dynamic Multi-cast -----
    n += 1
    slide_table(
        prs, n,
        title="③ Dynamic Multi-cast — 게임플레이의 표준",
        headers=["항목", "내용"],
        rows=[
            ["매크로", "DECLARE_DYNAMIC_MULTICAST_DELEGATE_TwoParams (...) 등"],
            ["바인딩", "여러 UFUNCTION만 가능 (AddDynamic) — 함수 이름 문자열로 찾음"],
            ["호출", ".Broadcast(파라미터들)"],
            ["BP 노출", "✅ — UPROPERTY(BlueprintAssignable) 로 BP에 이벤트 핀이 생긴다"],
            ["직렬화", "✅ — 저장/로드 가능"],
            ["대표 용도", "OnHealthChanged / OnDeath / OnInteract — 게임플레이 이벤트 대부분"],
        ],
        col_widths=[3.0, 9.83],
        star="우리가 만들 OnHealthChanged / OnDeath는 모두 이거. 이유는 — HUD가 BP니까",
    )

    # ----- Slide 16: 3종 비교표 -----
    n += 1
    slide_table(
        prs, n,
        title="델리게이트 3종 — 한 표로",
        headers=["종류", "매크로", "특징", "용도"],
        rows=[
            ["Single-cast", "DECLARE_DELEGATE", "1개 함수만, C++ 전용", "콜백 (비동기 로딩 완료)"],
            ["Multi-cast", "DECLARE_MULTICAST_DELEGATE", "N개 바인딩, C++ 전용", "엔진 내부 이벤트"],
            ["Dynamic\nMulti-cast", "DECLARE_DYNAMIC_MULTICAST_DELEGATE", "N개, BP 노출, 직렬화 가능, UFUNCTION만 바인딩 가능", "게임플레이 이벤트 (대부분)"],
        ],
        col_widths=[2.0, 4.5, 3.5, 3.33],
        star="기본 선택은 Dynamic Multi-cast. \"BP가 들을 일이 있나?\" → 있으면 무조건 Dynamic",
    )

    # ----- Slide 17: 왜 우리는 Dynamic Multi-cast? -----
    n += 1
    slide_callout(
        prs, n,
        title="우리는 왜 Dynamic Multi-cast를 쓰는가?",
        body_lines=[
            "•  HUD는 BP로 만든다 (UMG 위젯 = BP)",
            "•  BP가 \"HP 변했어!\" 신호를 들으려면 — 그 델리게이트가 BP에 노출되어 있어야 한다",
            "•  Single-cast / Multi-cast는 BP에 안 보인다 → BP 위젯이 구독 못함",
            "•  Dynamic Multi-cast는 BlueprintAssignable이 가능 → BP에서 \"이벤트 바인딩\" 노드로 구독",
            "",
            "•  비용? — 함수 이름을 문자열로 찾아 호출하므로 일반 호출보다 약간 느림. 게임플레이 이벤트 빈도(초당 수십 회)에선 무시할 수준.",
        ],
        metaphor="Dynamic Multi-cast = \"BP도 들을 수 있는 마이크\" — 약간의 비용을 내고 BP 호환성을 산다",
        star="기본은 Dynamic. \"성능이 진짜 문제될 때만\" Multi-cast 고려",
    )

    # ===== D — 이론 3: 컴포넌트 기반 설계 (slides 18–23) =====

    # ----- Slide 18: 상속 vs 조합 -----
    n += 1
    slide_two_col(
        prs, n,
        title="상속(is-a) vs 조합(has-a)",
        left_title="상속 — is-a",
        left_body=(
            "• \"A는 B의 일종이다\"\n"
            "• 부모 클래스를 골라야 함\n"
            "• 다중 상속 어려움 (C++)\n"
            "• 부모에 박힌 기능은 못 떼냄\n"
            "\n"
            "• 예: AStudyCharacter is-a ACharacter"
        ),
        right_title="조합 — has-a",
        right_body=(
            "• \"A는 B를 부품으로 가진다\"\n"
            "• 여러 컴포넌트를 동시에 끼움\n"
            "• 끼우고 빼는 게 자유로움\n"
            "• 같은 부품을 다른 액터도 사용\n"
            "\n"
            "• 예: AStudyCharacter has-a USpringArmComponent"
        ),
        star="언리얼은 이미 조합 기반 — Camera, SpringArm, Movement, ... 모두 컴포넌트",
    )

    # ----- Slide 19: 4주차 PlayerState의 HP -----
    n += 1
    slide_code(
        prs, n,
        title="4주차 PlayerState — HP가 박혀 있다",
        code=(
            "// StudyPlayerState.h\n"
            "UCLASS()\n"
            "class AStudyPlayerState : public APlayerState\n"
            "{\n"
            "    GENERATED_BODY()\n"
            "public:\n"
            "    UPROPERTY(EditAnywhere) float MaxHP = 100.f;\n"
            "    UPROPERTY(BlueprintReadOnly) float CurrentHP = 100.f;\n"
            "\n"
            "    UFUNCTION(BlueprintCallable) void ApplyDamage(float Amount);\n"
            "    UFUNCTION(BlueprintCallable) void Heal(float Amount);\n"
            "    UFUNCTION(BlueprintPure)     bool  IsAlive() const;\n"
            "    UFUNCTION(BlueprintCallable) void  ResetHP();\n"
            "};"
        ),
        comment="잘 작동한다. 그런데 — 이 로직, 다음 주 \"적\"도 똑같이 필요하지 않나?",
        star="HP는 PlayerState의 본질적 특성이 아니다. \"체력이 있는 모든 것\"의 공통 기능이다",
    )

    # ----- Slide 20: 적은 PlayerState가 없다 -----
    n += 1
    slide_quote(
        prs, n,
        title="다음 주 질문 — 적은 PlayerState가 없다",
        quote="적도 HP가 있다. 그런데 적은 PlayerState가 없는데?",
        subtext="PlayerState는 \"플레이어 1명당 1개\" 만들어진다.\n"
                "AI 적 100마리에게 PlayerState를 100개 만들 순 없다 — PlayerController가 없으니까.\n\n"
                "그렇다고 적용 HP 로직을 또 만들면 — 같은 코드가 두 곳에 존재한다 (DRY 위반).\n"
                "해결: HP 로직을 \"부착 가능한 부품\"으로 떼어내자.",
    )

    # ----- Slide 21: UActorComponent의 약속 -----
    n += 1
    slide_bullets(
        prs, n,
        title="UActorComponent — \"부착 가능한 기능 단위\"",
        bullets=[
            "어떤 AActor에든 부착 가능 — Character, PlayerState, Pawn, 무엇이든",
            "자기 라이프사이클 — BeginPlay / EndPlay / TickComponent 가짐",
            "자기 데이터 — UPROPERTY 가지고, 디테일 패널에 자기 섹션이 뜸",
            "자기 메서드 — BlueprintCallable 함수를 BP에서 호출 가능",
            "★ 핵심 — Owner를 모른다. \"체력이 있는 액터에 붙는다\"고만 알고, Player든 Enemy든 신경 안 씀",
        ],
        star="컴포넌트는 \"누구의 부품\"인지 모르는 게 미덕 — 그래서 어디든 붙는다",
    )

    # ----- Slide 22: 4개 액터에 붙이기 -----
    n += 1
    slide_four_box(
        prs, n,
        title="같은 HealthComponent를 4종 액터에 붙인다",
        boxes=[
            ("Player", "AStudyPlayerState\n→ Health 부착\n→ 데미지/회복/죽음 흐름은 동일\n→ HUD가 OnHealthChanged 구독"),
            ("Enemy", "AEnemyCharacter (8주차)\n→ 같은 Health 부착\n→ 적도 같은 데미지 시스템 사용\n→ GameState가 OnDeath 구독 → 처치 수 ++"),
            ("Boss", "ABossCharacter\n→ 같은 Health 부착 + MaxHP=5000\n→ 보스 HP UI도 OnHealthChanged 구독\n→ 데이터만 다르고 코드는 같음"),
            ("Destructible", "ADestructibleCrate\n→ Pawn도 아닌 단순 Actor\n→ 같은 Health 부착 → 부수면 사라짐\n→ \"체력 있는 모든 것\"에 일관된 시스템"),
        ],
        star="추출의 핵심 — 부착 가능한 부품이 되면, \"체력이 필요한 모든 것\"에 한 번에 답이 된다",
    )

    # ----- Slide 23: 추출의 핵심 -----
    n += 1
    slide_callout(
        prs, n,
        title="추출의 핵심 — 상속이 아니라 \"부착\"",
        body_lines=[
            "•  만약 \"HP 있는 액터\"라는 부모 클래스를 만들면? → AHealthActor를 PlayerState도 상속, Enemy도 상속...",
            "•  → 단일 상속 제약, 클래스 트리가 뚱뚱해짐, 새 기능마다 새 부모",
            "•  컴포넌트는 그 반대 — \"기능을 끼우고 뺀다\"",
            "•  PlayerState는 여전히 APlayerState이고, Enemy는 여전히 ACharacter다. 그저 둘 다 Health를 끼울 뿐",
            "",
            "•  ★ 좋은 추상화는 \"무엇이다\"가 아니라 \"무엇을 한다\"에 집중한다",
        ],
        metaphor="레고 블록 — 한 가지 색만 쓸 수도, 다섯 가지 색을 동시에 끼울 수도 있다",
    )

    # ===== E — 분기 (slide 24) =====

    # ----- Slide 24: Divider — 실습 시작 -----
    n += 1
    slide_divider(
        prs, n,
        title="실습 — 직접 떼어내자",
        subtitle="60분 (HealthComponent 만들기 → PlayerState 이전 → HUD 구독 → 포션 연결)",
    )

    # ===== F — 실습 1: UHealthComponent 추출 (slides 25–31) =====

    # ----- Slide 25: Step 1 — 클래스 생성 -----
    n += 1
    slide_bullets(
        prs, n,
        title="실습 1-① — UActorComponent 클래스 생성 (3분)",
        bullets=[
            "Tools → New C++ Class…",
            "부모: Actor Component  (Show All Classes 필요할 수 있음)",
            "이름: HealthComponent  (앞에 U는 자동으로 붙음)",
            "위치: Public",
            "★ 컴파일 한 번 — Editor 재시작 안 해도 됨 (Live Coding 또는 Hot Reload)",
        ],
        sub_top=1.6,
        body_size=22,
        star="Wizard가 생성한 코드는 거의 안 만진다 — 핵심은 헤더와 cpp의 본문",
    )

    # ----- Slide 26: 헤더 — 멤버변수 + 시그니처 -----
    n += 1
    slide_code(
        prs, n,
        title="실습 1-② — HealthComponent.h: 멤버 + 시그니처",
        code=(
            "UCLASS(ClassGroup=(Custom), meta=(BlueprintSpawnableComponent))\n"
            "class STUDY_API UHealthComponent : public UActorComponent\n"
            "{\n"
            "    GENERATED_BODY()\n"
            "public:\n"
            "    UHealthComponent();\n"
            "\n"
            "    UPROPERTY(EditAnywhere, BlueprintReadOnly, Category=\"Health\")\n"
            "    float MaxHP = 100.f;\n"
            "\n"
            "    UPROPERTY(VisibleAnywhere, BlueprintReadOnly, Category=\"Health\")\n"
            "    float CurrentHP = 100.f;\n"
            "\n"
            "    UFUNCTION(BlueprintCallable) void ApplyDamage(float Amount);\n"
            "    UFUNCTION(BlueprintCallable) void Heal(float Amount);\n"
            "    UFUNCTION(BlueprintPure)     bool  IsAlive() const;\n"
            "    UFUNCTION(BlueprintCallable) void  ResetHP();\n"
            "};"
        ),
        comment="meta=(BlueprintSpawnableComponent) — 이거 있어야 BP \"Add Component\" 메뉴에 뜬다",
        code_size=14,
    )

    # ----- Slide 27: 헤더 — OnHealthChanged 델리게이트 -----
    n += 1
    slide_code(
        prs, n,
        title="실습 1-③ — OnHealthChanged 델리게이트 선언",
        code=(
            "// HealthComponent.h — 클래스 \"바깥\"에 매크로 선언\n"
            "DECLARE_DYNAMIC_MULTICAST_DELEGATE_TwoParams(\n"
            "    FOnHealthChanged,\n"
            "    float, NewHP,\n"
            "    float, MaxHPParam\n"
            ");\n"
            "\n"
            "UCLASS(...)\n"
            "class STUDY_API UHealthComponent : public UActorComponent\n"
            "{\n"
            "    GENERATED_BODY()\n"
            "public:\n"
            "    // ... 기존 멤버 ...\n"
            "\n"
            "    UPROPERTY(BlueprintAssignable, Category=\"Health\")\n"
            "    FOnHealthChanged OnHealthChanged;\n"
            "};"
        ),
        comment="DECLARE_DYNAMIC_MULTICAST_DELEGATE_TwoParams — 파라미터 2개 (각각 \"타입, 이름\" 쌍)",
        star="BlueprintAssignable이 있으면 BP에서 \"Bind Event to On Health Changed\" 노드가 생긴다",
        code_size=14,
    )

    # ----- Slide 28: 헤더 — OnDeath 델리게이트 -----
    n += 1
    slide_code(
        prs, n,
        title="실습 1-④ — OnDeath 델리게이트 (파라미터 없음)",
        code=(
            "// HealthComponent.h — 또 하나 추가\n"
            "DECLARE_DYNAMIC_MULTICAST_DELEGATE(FOnDeath);\n"
            "\n"
            "UCLASS(...)\n"
            "class STUDY_API UHealthComponent : public UActorComponent\n"
            "{\n"
            "    GENERATED_BODY()\n"
            "public:\n"
            "    // ... OnHealthChanged ...\n"
            "\n"
            "    UPROPERTY(BlueprintAssignable, Category=\"Health\")\n"
            "    FOnDeath OnDeath;\n"
            "};\n"
            "\n"
            "// 파라미터가 없으면 매크로 끝의 _OneParam, _TwoParams 같은 접미사 안 붙임"
        ),
        comment="\"파라미터 N개\"에 따라 매크로 이름이 다르다 — _OneParam, _TwoParams, _ThreeParams ... ",
        code_size=14,
    )

    # ----- Slide 29: cpp — ApplyDamage Broadcast -----
    n += 1
    slide_code(
        prs, n,
        title="실습 1-⑤ — HealthComponent.cpp: ApplyDamage",
        code=(
            "void UHealthComponent::ApplyDamage(float Amount)\n"
            "{\n"
            "    if (!IsAlive()) return;  // 이미 죽었으면 무시\n"
            "\n"
            "    CurrentHP = FMath::Clamp(CurrentHP - Amount, 0.f, MaxHP);\n"
            "\n"
            "    // \"변했다!\" — 모든 구독자에게 한 번에\n"
            "    OnHealthChanged.Broadcast(CurrentHP, MaxHP);\n"
            "\n"
            "    if (CurrentHP <= 0.f)\n"
            "    {\n"
            "        OnDeath.Broadcast();\n"
            "    }\n"
            "}"
        ),
        comment=".Broadcast(...) — 등록된 모든 함수가 호출된다. 호출 순서는 \"보장 안 됨\"",
        star="값을 바꾼 \"직후\"에 Broadcast — \"이 시점부터 새 값이다\"라는 약속",
    )

    # ----- Slide 30: cpp — Heal / ResetHP -----
    n += 1
    slide_code(
        prs, n,
        title="실습 1-⑥ — Heal / IsAlive / ResetHP",
        code=(
            "void UHealthComponent::Heal(float Amount)\n"
            "{\n"
            "    if (!IsAlive()) return;  // 죽은 후 회복 못함 (디자인 결정)\n"
            "    CurrentHP = FMath::Clamp(CurrentHP + Amount, 0.f, MaxHP);\n"
            "    OnHealthChanged.Broadcast(CurrentHP, MaxHP);\n"
            "}\n"
            "\n"
            "bool UHealthComponent::IsAlive() const\n"
            "{\n"
            "    return CurrentHP > 0.f;\n"
            "}\n"
            "\n"
            "void UHealthComponent::ResetHP()\n"
            "{\n"
            "    CurrentHP = MaxHP;\n"
            "    OnHealthChanged.Broadcast(CurrentHP, MaxHP);\n"
            "}"
        ),
        comment="HP를 변경하는 모든 메서드는 마지막에 OnHealthChanged.Broadcast — 일관성 유지",
    )

    # ----- Slide 31: 컴파일 확인 -----
    n += 1
    slide_callout(
        prs, n,
        title="실습 1-⑦ — 컴파일 → 컴포넌트 등록 확인",
        body_lines=[
            "•  Live Coding 또는 Build (Ctrl+Alt+F11) — 에러 없이 빌드되어야 함",
            "•  BP_ThirdPersonCharacter 열기 → Add Component → \"Health\" 검색",
            "•  → 목록에 \"Health\" 가 떠야 함  (meta=(BlueprintSpawnableComponent) 효과)",
            "",
            "•  아직 PlayerState에는 안 붙였다 — 그건 실습 2",
        ],
        metaphor="여기까지가 \"부품을 만든 단계\". 다음은 \"부품을 끼우는 단계\"",
        star="\"검색에 뜬다\" = 리플렉션 등록 성공 + 매크로 정상 = OK 신호",
    )

    # ===== G — 실습 2: PlayerState에서 이전 (slides 32–35) =====

    # ----- Slide 32: PS.h — UPROPERTY 추가 -----
    n += 1
    slide_code(
        prs, n,
        title="실습 2-① — StudyPlayerState.h에 Health 추가",
        code=(
            "// StudyPlayerState.h\n"
            "#include \"StudyPlayerState.generated.h\"\n"
            "\n"
            "class UHealthComponent;  // forward declare\n"
            "\n"
            "UCLASS()\n"
            "class AStudyPlayerState : public APlayerState\n"
            "{\n"
            "    GENERATED_BODY()\n"
            "public:\n"
            "    AStudyPlayerState();\n"
            "\n"
            "    UPROPERTY(VisibleAnywhere, BlueprintReadOnly, Category=\"Stats\")\n"
            "    UHealthComponent* Health;  // ← 추가\n"
            "\n"
            "    // 기존 ApplyDamage/Heal/IsAlive/ResetHP 시그니처는 그대로 둔다\n"
            "    UFUNCTION(BlueprintCallable) void ApplyDamage(float Amount);\n"
            "    // ...\n"
            "};"
        ),
        comment="forward declare로 헤더 의존성 최소화. .cpp에서 #include \"HealthComponent.h\"",
    )

    # ----- Slide 33: PS.cpp — CreateDefaultSubobject -----
    n += 1
    slide_code(
        prs, n,
        title="실습 2-② — 생성자에서 컴포넌트 부착",
        code=(
            "// StudyPlayerState.cpp\n"
            "#include \"StudyPlayerState.h\"\n"
            "#include \"HealthComponent.h\"  // ← 여기서 진짜 include\n"
            "\n"
            "AStudyPlayerState::AStudyPlayerState()\n"
            "{\n"
            "    Health = CreateDefaultSubobject<UHealthComponent>(TEXT(\"Health\"));\n"
            "}\n"
            "\n"
            "// CreateDefaultSubobject — 생성자에서만 호출 가능.\n"
            "// 이 시점에 컴포넌트가 PlayerState에 \"부착\"된다.\n"
            "// 이름 \"Health\"는 디테일 패널에 표시됨."
        ),
        comment="CreateDefaultSubobject<T>(TEXT(\"이름\")) — UE의 \"부품 만들고 끼우기\" 정석",
        star="이게 끝이다 — PlayerState가 Health를 \"가지게\" 됨. 다음은 기존 메서드를 위임만 하면 됨",
    )

    # ----- Slide 34: PS.cpp — wrapper로 -----
    n += 1
    slide_code(
        prs, n,
        title="실습 2-③ — 기존 메서드를 1줄 wrapper로",
        code=(
            "// StudyPlayerState.cpp — 기존 ApplyDamage 본문을 이렇게 바꾼다\n"
            "void AStudyPlayerState::ApplyDamage(float Amount)\n"
            "{\n"
            "    if (Health) Health->ApplyDamage(Amount);\n"
            "}\n"
            "\n"
            "void AStudyPlayerState::Heal(float Amount)\n"
            "{\n"
            "    if (Health) Health->Heal(Amount);\n"
            "}\n"
            "\n"
            "bool AStudyPlayerState::IsAlive() const\n"
            "{\n"
            "    return Health && Health->IsAlive();\n"
            "}\n"
            "\n"
            "void AStudyPlayerState::ResetHP()\n"
            "{\n"
            "    if (Health) Health->ResetHP();\n"
            "}"
        ),
        comment="기존 호출자(Character::OnTickDamage 등)는 한 줄도 안 바꾼다 — wrapper 덕분",
    )

    # ----- Slide 35: 위임의 미덕 -----
    n += 1
    slide_callout(
        prs, n,
        title="실습 2-④ — Character::OnTickDamage 코드는 그대로",
        body_lines=[
            "•  4주차의 PS->ApplyDamage(10) 호출 — 한 줄도 안 바뀐다",
            "•  내부적으로는 Health->ApplyDamage가 불리고, OnHealthChanged.Broadcast가 자동",
            "•  IsAlive()도 그대로 작동 → HandlePlayerDeath() 흐름 그대로",
            "",
            "•  ★ \"추출\"은 \"파괴\"가 아니라 \"위임\". 호출자에게는 보이지 않게 책임만 옮긴다.",
        ],
        metaphor="회사에서 부서를 신설한 것 같다 — 외부 거래처는 여전히 같은 번호로 전화하고, 안에서 새 부서가 처리",
        star="이게 \"호환성을 깨지 않는 리팩토링\"의 정석. 큰 변화도 작은 영향으로 만들 수 있다",
    )

    # ===== H — 실습 3: HUD가 델리게이트 구독 (slides 36–41) =====

    # ----- Slide 36: Construct 흐름 -----
    n += 1
    slide_flow(
        prs, n,
        title="실습 3-① — WBP_HpBar 의 Construct 이벤트",
        flow_lines=[
            "Event Construct  (위젯이 화면에 추가될 때 1번)",
            "↓",
            "Get Owning Player  →  Get Player State",
            "↓",
            "Cast To StudyPlayerState",
            "↓",
            "Health (UHealthComponent*)",
            "↓",
            "OnHealthChanged  →  Bind Event to ...  →  HandleHealthChanged",
        ],
        star="구독은 \"한 번\" 한다 — Construct가 자연스러운 시점",
    )

    # ----- Slide 37: AddDynamic 코드 -----
    n += 1
    slide_code(
        prs, n,
        title="실습 3-② — BP에서 \"Bind Event\" 노드 (의사코드)",
        code=(
            "// 개념적으로 C++로 쓰면 이런 구조다 (BP에서는 노드로 연결)\n"
            "void UWBP_HpBar::NativeConstruct()\n"
            "{\n"
            "    Super::NativeConstruct();\n"
            "\n"
            "    APlayerController* PC = GetOwningPlayer();\n"
            "    AStudyPlayerState* PS = PC ? PC->GetPlayerState<AStudyPlayerState>() : nullptr;\n"
            "    if (PS && PS->Health)\n"
            "    {\n"
            "        // ★ AddDynamic — UFUNCTION만 받는다. 이름을 문자열로 찾음\n"
            "        PS->Health->OnHealthChanged.AddDynamic(\n"
            "            this, &UWBP_HpBar::HandleHealthChanged);\n"
            "    }\n"
            "}"
        ),
        comment="BP 그래프에서는 OnHealthChanged 핀에서 드래그 → \"Bind Event to ...\" → 이벤트 노드 생성",
        star="HandleHealthChanged 는 UFUNCTION이어야 함 (BP의 \"Custom Event\"가 자동으로 만족)",
    )

    # ----- Slide 38: HandleHealthChanged 본문 -----
    n += 1
    slide_code(
        prs, n,
        title="실습 3-③ — HandleHealthChanged: ProgressBar.SetPercent",
        code=(
            "// BP 의사코드 — \"On Health Changed\" 이벤트 노드의 본문\n"
            "// 파라미터: NewHP (float), MaxHPParam (float)\n"
            "//\n"
            "// 노드 연결:\n"
            "//   New HP  ──┐\n"
            "//             ├──▶  Divide  ──▶  Set Percent (HP_ProgressBar)\n"
            "//   Max HP  ──┘\n"
            "//\n"
            "// 또는 한 줄로:\n"
            "//   HP_ProgressBar.SetPercent(NewHP / MaxHPParam)\n"
            "\n"
            "// 이게 \"실제 HP가 변했을 때만\" 호출된다.\n"
            "// 매 프레임 폴링 ❌  →  변경 시점만 ✅"
        ),
        comment="ProgressBar 위젯의 \"is variable\" 체크 → BP에서 변수처럼 접근 가능",
    )

    # ----- Slide 39: Percent Binding 제거 -----
    n += 1
    slide_bullets(
        prs, n,
        title="실습 3-④ — 기존 Percent Binding 제거",
        bullets=[
            "WBP_HpBar의 ProgressBar 선택",
            "Details → Progress → Percent → 우클릭 → \"Remove Binding\"",
            "이걸 안 빼면 — 둘 다 작동해서 매 프레임 호출이 그대로 남는다",
            "이제 Percent는 \"기본값 1.0\"이고, HandleHealthChanged가 호출될 때만 SetPercent로 변경됨",
            "★ Compile → Save → PIE",
        ],
        sub_top=1.6,
        body_size=22,
        star="\"무엇을 빼는가\"가 \"무엇을 더하는가\"만큼 중요하다",
    )

    # ----- Slide 40: 실험 결과 — 호출 횟수 비교 -----
    n += 1
    slide_experiment_result(
        prs, n,
        title="실험 결과 — 슬라이드 9와 비교",
        observation="HandleHealthChanged 안에 Print String 추가 → 1초에 1번만 찍힌다 (1초마다 데미지 1회)",
        meaning="\"실제로 변할 때만\" 호출된다. 폴링 60회/초 → 푸시 1회/초 — 60배 절약.",
        star="푸시는 \"이벤트 발생 횟수\"만큼만 비싸다 — 안 변하는 동안은 0번",
    )

    # ----- Slide 41: Dangling delegate 예고 -----
    n += 1
    slide_callout(
        prs, n,
        title="실습 3-⑤ — Destruct에서 RemoveDynamic 안 하면?",
        body_lines=[
            "•  HUD 위젯이 사라질 때 (RemoveFromParent 등) — 구독은 자동으로 끊어지지 않음",
            "•  PS->Health는 살아있고, OnHealthChanged 안의 \"this 포인터\"는 죽은 위젯을 가리킬 수 있음",
            "•  → 그러나 Dynamic Multi-cast는 \"WeakObjectPtr\" 기반이라 GC가 정리한 객체는 자동으로 호출 안 함 (UE 보호장치)",
            "•  ★ 그래도 RemoveDynamic을 명시적으로 호출하는 게 좋은 습관 — 구독자가 \"왜 살아있는지\" 확인 가능",
        ],
        metaphor="유튜브 채널 구독을 안 끊으면 알림은 안 와도 \"구독자 수\"엔 잡혀있다 — 일종의 누수",
        star="이건 \"핵심 질문 1\" — 슬라이드 46에서 다시 다룬다",
    )

    # ===== I — 실습 4: 포션 연결 (slides 42–44) =====

    # ----- Slide 42: Week 6 backfill 안내 -----
    n += 1
    slide_callout(
        prs, n,
        title="실습 4-① — ⚠️ 사전 작업: BP_HealthPotion 만들기",
        body_lines=[
            "•  6주차에 만들어야 했던 BP_HealthPotion이 아직 Content/에 없다",
            "•  먼저 만들어 두자 — 6주차 \"베이스 C++ + 변종 BP\" 패턴 그대로",
            "",
            "•  Content Browser → 우클릭 → Blueprint Class → APickupItem 검색 → 선택",
            "•  이름: BP_HealthPotion",
            "•  안에 변수 HealAmount (float, default 30) 추가",
            "•  Mesh도 적당한 거 골라 부착 (큐브, 파편, 뭐든)",
        ],
        metaphor="6주차 미완 숙제 5분 → 7주차 실습 30초",
        star="이게 \"베이스 C++ + 변종 BP\" 패턴의 효율 — 한 번 만든 베이스가 매주 재활용된다",
    )

    # ----- Slide 43: OnPickedUp override -----
    n += 1
    slide_code(
        prs, n,
        title="실습 4-② — BP_HealthPotion: OnPickedUp 오버라이드",
        code=(
            "// BP 의사코드 — BP_HealthPotion 의 OnPickedUp (오버라이드)\n"
            "// 파라미터: AActor* Picker\n"
            "\n"
            "// 노드 연결:\n"
            "Cast To Pawn (Picker)\n"
            "    └─▶ Get Player State\n"
            "         └─▶ Cast To StudyPlayerState\n"
            "              └─▶ Health  (UHealthComponent*)\n"
            "                   └─▶ Heal (HealAmount)\n"
            "\n"
            "// 마지막에:\n"
            "Parent: OnPickedUp   // C++의 OnPickedUp_Implementation 호출 (자기파괴 + 효과)\n"
            "\n"
            "// 결과: HP가 회복되고, OnHealthChanged.Broadcast 자동 발동\n"
            "//       → HUD 자동 갱신 (실습 3에서 구독했으니까)"
        ),
        comment="포션은 PlayerState의 CurrentHP를 직접 만지지 않는다. Health->Heal만 안다",
        star="\"누가 누구를 아는지\"가 정리됨 — 포션은 PlayerState 타입을 알지만, HP 로직은 안 만진다",
    )

    # ----- Slide 44: Before / After 결합도 -----
    n += 1
    slide_two_col(
        prs, n,
        title="결합도 비교 — Before / After",
        left_title="Before (4주차 식)",
        left_body=(
            "• 포션이 PlayerState->Heal() 직접 호출\n"
            "• \"HP 회복\"이 PS의 책임\n"
            "• 적은 PS가 없으니 같은 포션 못 씀\n"
            "• 새 \"체력 있는 액터\"마다 코드 복사\n"
            "\n"
            "• 결합 그래프:\n"
            "  포션 ──▶ PlayerState ──▶ HP 로직"
        ),
        right_title="After (이번 주)",
        right_body=(
            "• 포션이 HealthComponent->Heal() 호출\n"
            "• \"HP 회복\"이 Health의 책임\n"
            "• Enemy도 Health 부착 → 같은 포션 OK\n"
            "• 새 \"체력 있는 액터\"는 컴포넌트만 추가\n"
            "\n"
            "• 결합 그래프:\n"
            "  포션 ──▶ HealthComponent ◀── Player\n"
            "                            ◀── Enemy"
        ),
        star="\"누가 누구를 아는지\"의 그래프가 단순해진다 — 좋은 설계의 신호",
    )

    # ===== J — 보너스 + 핵심 질문 (slides 45–48) =====

    # ----- Slide 45: 보너스 — 트리거 볼륨 -----
    n += 1
    slide_flow(
        prs, n,
        title="(보너스) 트리거 볼륨 — 용암 함정 만들기",
        flow_lines=[
            "레벨에 BP Actor 배치 (BoxComponent 포함)",
            "↓",
            "BoxComponent.OnComponentBeginOverlap  →  이벤트 바인딩",
            "↓",
            "Other Actor  →  Cast To Character  →  Get Player State",
            "↓",
            "Cast To StudyPlayerState  →  Health  →  ApplyDamage(50)",
            "↓",
            "→ HP 50 감소 + HUD 자동 갱신 + (HP 0이면) OnDeath.Broadcast",
        ],
        star="BoxComponent의 OnComponentBeginOverlap도 \"Dynamic Multi-cast 델리게이트\" — 같은 패턴",
    )

    # ----- Slide 46: 핵심 질문 1 — Dangling -----
    n += 1
    slide_quote(
        prs, n,
        title="핵심 질문 ① — Dangling Delegate",
        quote="HUD 위젯이 사라질 때 RemoveDynamic을 안 하면 어떻게 되는가?",
        subtext="UE는 Dynamic Multi-cast의 \"this\" 포인터를 WeakObjectPtr로 보관한다.\n"
                "GC가 위젯을 정리하면 — 다음 Broadcast 때 자동으로 \"죽은 구독자\"는 호출에서 제외된다.\n\n"
                "즉, 크래시는 안 난다. 하지만 — 명시적 RemoveDynamic이 \"누가 살아있어야 하는지\"를 코드에 박아둔다.\n"
                "★ 좋은 코드는 \"동작하는 것\"보다 \"읽힐 때 의도가 보이는 것\"이다.",
    )

    # ----- Slide 47: 핵심 질문 2 — GameState 구독 -----
    n += 1
    slide_quote(
        prs, n,
        title="핵심 질문 ② — 다음 주 매니저-컴포넌트 협력",
        quote="적의 OnDeath를 GameState가 구독하면 어떻게 되는가?",
        subtext="EnemyCharacter::Health->OnDeath.AddDynamic(GameState, &AStudyGameState::OnEnemyDied)\n"
                "→ GameState::EnemiesRemaining-- → 0이면 \"승리\" 이벤트\n\n"
                "5주차의 \"6 매니저\"가 7주차의 \"컴포넌트\"와 만나는 순간 —\n"
                "매니저는 게임 \"전체\"를 보고, 컴포넌트는 액터 \"하나\"를 본다. 둘이 협력하면 게임 루프가 닫힌다.",
    )

    # ----- Slide 48: 핵심 질문 3 — 호출 순서 -----
    n += 1
    slide_quote(
        prs, n,
        title="핵심 질문 ③ — Multicast 호출 순서",
        quote="Multicast가 N개 함수를 호출할 때 — 순서가 보장되는가?",
        subtext="결론: 보장 안 됨. 등록 순서대로 호출되긴 하지만, 그건 \"구현 디테일\"이지 \"계약\"이 아님.\n\n"
                "★ 따라서 — 한 구독자가 다른 구독자보다 \"먼저\" 호출될 거라 가정하지 말 것.\n"
                "구독자끼리 의존성이 필요하면, 델리게이트가 아니라 명시적 호출 체인을 만들어라.",
    )

    # ===== K — 마무리 (slides 49–50) =====

    # ----- Slide 49: 다음 주 예고 -----
    n += 1
    slide_callout(
        prs, n,
        title="다음 주 예고 — 8주차 (5/25)",
        body_lines=[
            "•  HealthComponent의 MaxHP=100이 코드에 박혀 있다 → DataAsset으로 분리",
            "•  → 디자이너가 코드 안 만지고 \"보스용 Health 5000\" 같은 변종을 만들 수 있게",
            "•  Enhanced Input 심화 — 모드별 컨트롤 (전투 / 메뉴 / 인벤토리) 전환",
            "•  적이 등장한다 — 같은 HealthComponent 부착, GameState가 OnDeath 구독해서 처치 카운트",
        ],
        metaphor="\"하드코딩된 값을 데이터로\" — 컴파일 없이 게임을 튜닝할 수 있는 첫 단계",
        star="6주차에 베이스를 만들었고, 7주차에 컴포넌트로 추출했다 — 8주차는 데이터로 외부화",
    )

    # ----- Slide 50: 정리 & Q&A -----
    n += 1
    s = blank_slide(prs)
    add_textbox(s, Inches(0.5), Inches(0.8), Inches(12.33), Inches(1.2),
                "오늘 가져갈 한 줄 — \"변할 때만 알린다\"",
                size=40, bold=True, color=COLOR_TITLE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(1.0), Inches(2.4), Inches(11.33), Inches(2.8),
                "✓ 폴링 → 푸시  (매 프레임 묻기 → 변할 때만 알림)\n"
                "✓ 상속 → 조합  (\"무엇이다\" → \"무엇을 한다\")\n"
                "✓ Dynamic Multi-cast — BP가 들을 수 있는 마이크\n"
                "✓ HealthComponent — 다음 주 적도 그대로 쓴다",
                size=24, color=COLOR_BODY, align=PP_ALIGN.LEFT)
    add_textbox(s, Inches(0.5), Inches(5.6), Inches(12.33), Inches(1.2),
                "정리 & Q&A",
                size=44, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(6.5), Inches(12.33), Inches(0.6),
                "다음 주 8주차: 5/25 (월) 18:30 — DataAsset · Enhanced Input 심화 · 적 등장",
                size=16, italic=True, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER)
    add_page_number(s, n)

    return n


def main():
    prs = Presentation(TEMPLATE)
    print(f"Loaded template: {len(list(prs.slides))} existing slides")
    remove_all_slides(prs)
    print(f"After strip: {len(list(prs.slides))} slides")
    total = build_slides(prs)
    print(f"Built {total} slides")
    prs.save(OUTPUT)
    print(f"Saved to: {OUTPUT}")


if __name__ == "__main__":
    main()
