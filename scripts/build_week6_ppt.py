"""Build 언리얼_스터디_6주차_자료.pptx from 5주차 template (16:9, 13.33x7.5in).

Strategy:
- Open 5주차 .pptx as template (theme/master inherited)
- Strip all existing slides
- Add 50 new slides per the slide-by-slide spec in the plan
- All slides use the single 'DEFAULT' layout, with text boxes positioned manually
"""
import io
import sys
from copy import deepcopy

# Force UTF-8 stdout for Windows cp949 console
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

TEMPLATE = r"C:\Users\kk409\Documents\Konkuk\3-1\Unreal_Study\언리얼_스터디_5주차_자료_수정본.pptx"
OUTPUT = r"C:\Users\kk409\Documents\Konkuk\3-1\Unreal_Study\언리얼_스터디_6주차_자료.pptx"

# Fonts
FONT_KR = "맑은 고딕"
FONT_CODE = "Consolas"

# Colors
COLOR_TITLE = RGBColor(0x1F, 0x1F, 0x1F)        # near-black for titles
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)         # body text
COLOR_ACCENT = RGBColor(0x2B, 0x57, 0x9A)       # blue accent
COLOR_STAR = RGBColor(0xC0, 0x39, 0x2B)         # red-ish for ★ insight
COLOR_CHECK = RGBColor(0x2E, 0x7D, 0x32)        # green for ✓
COLOR_CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)      # light gray for code blocks
COLOR_TABLE_HEADER_BG = RGBColor(0x2B, 0x57, 0x9A)
COLOR_TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TABLE_ROW_ALT = RGBColor(0xF0, 0xF4, 0xF8)
COLOR_DIVIDER_BG = RGBColor(0x1F, 0x1F, 0x1F)
COLOR_DIVIDER_FG = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_QUOTE_BG = RGBColor(0xFF, 0xF8, 0xE1)

# Slide region constants (16:9, 13.33x7.5in)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.5)


def remove_all_slides(prs):
    """Remove every slide so we start clean."""
    sldIdLst = prs.slides._sldIdLst
    rId_list = [sld.attrib[qn('r:id')] for sld in list(sldIdLst)]
    for rId in rId_list:
        prs.part.drop_rel(rId)
    for sld in list(sldIdLst):
        sldIdLst.remove(sld)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[0])


def set_para_font(para, name=FONT_KR, size=18, bold=False, italic=False, color=COLOR_BODY):
    for run in para.runs:
        run.font.name = name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        # Also set East Asian font in rPr
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
            rPr.append(ea)
        else:
            ea.set('typeface', FONT_KR)


def add_textbox(slide, left, top, width, height, text, *, size=18, bold=False, italic=False,
                color=COLOR_BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT_KR,
                fill=None, line_color=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if line_color is not None:
        tb.line.color.rgb = line_color
    else:
        tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        set_para_font(p, name=font, size=size, bold=bold, italic=italic, color=color)
    return tb


def add_title(slide, text, *, size=32, color=COLOR_TITLE):
    return add_textbox(
        slide, Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.85),
        text, size=size, bold=True, color=color, align=PP_ALIGN.LEFT,
    )


def add_underline(slide):
    """Thin accent line under the title."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.18), Inches(12.33), Emu(20000))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()


def add_star_footer(slide, text):
    add_textbox(
        slide, Inches(0.5), Inches(6.7), Inches(12.33), Inches(0.6),
        f"★  {text}", size=16, italic=True, bold=True, color=COLOR_STAR,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
        fill=RGBColor(0xFF, 0xF4, 0xE0),
    )


def add_page_number(slide, n, total=50):
    add_textbox(
        slide, Inches(12.33), Inches(7.05), Inches(0.95), Inches(0.35),
        f"{n} / {total}", size=10, color=RGBColor(0x99, 0x99, 0x99),
        align=PP_ALIGN.RIGHT,
    )


# ----------- Slide builders -----------

def slide_title_page(prs, title, subtitle, footer):
    s = blank_slide(prs)
    add_textbox(s, Inches(0.5), Inches(2.4), Inches(12.33), Inches(1.4),
                title, size=54, bold=True, color=COLOR_TITLE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(3.9), Inches(12.33), Inches(0.9),
                subtitle, size=28, italic=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(5.5), Inches(12.33), Inches(0.6),
                footer, size=18, color=RGBColor(0x66, 0x66, 0x66), align=PP_ALIGN.CENTER, font=FONT_CODE)
    add_textbox(s, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0.4),
                "박준건  |  건국대 언리얼 스터디  |  2026-05-14", size=12,
                color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)
    return s


def slide_divider(prs, n, title, subtitle):
    s = blank_slide(prs)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_DIVIDER_BG
    bg.line.fill.background()
    add_textbox(s, Inches(0.5), Inches(2.6), Inches(12.33), Inches(1.6),
                title, size=60, bold=True, color=COLOR_DIVIDER_FG, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(4.4), Inches(12.33), Inches(0.8),
                subtitle, size=24, italic=True, color=RGBColor(0xCC, 0xCC, 0xCC), align=PP_ALIGN.CENTER)
    return s


def slide_quote(prs, n, title, quote, subtext):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    # Quote box
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.0), Inches(2.3), Inches(11.33), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_QUOTE_BG
    box.line.color.rgb = COLOR_ACCENT
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f'"{quote}"'
    set_para_font(p, size=36, bold=True, color=COLOR_TITLE)
    add_textbox(s, Inches(1.0), Inches(4.7), Inches(11.33), Inches(2.0),
                subtext, size=18, italic=True, color=COLOR_BODY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    add_page_number(s, n)
    return s


def slide_bullets(prs, n, title, bullets, *, star=None, sub_top=1.5, body_size=20):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    body_top = Inches(sub_top)
    body_h = Inches(7.0 - sub_top - (0.7 if star else 0.2))
    tb = s.shapes.add_textbox(Inches(0.7), body_top, Inches(12.0), body_h)
    tb.fill.background()
    tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # parse leading marker
        text = b
        marker_color = COLOR_BODY
        if text.startswith("✓ "):
            run = p.add_run()
            run.text = "✓  "
            run.font.name = FONT_KR
            run.font.size = Pt(body_size)
            run.font.bold = True
            run.font.color.rgb = COLOR_CHECK
            text = text[2:]
        elif text.startswith("★ "):
            run = p.add_run()
            run.text = "★  "
            run.font.name = FONT_KR
            run.font.size = Pt(body_size)
            run.font.bold = True
            run.font.color.rgb = COLOR_STAR
            text = text[2:]
        else:
            run = p.add_run()
            run.text = "•  "
            run.font.name = FONT_KR
            run.font.size = Pt(body_size)
            run.font.color.rgb = COLOR_ACCENT
        run = p.add_run()
        run.text = text
        run.font.name = FONT_KR
        run.font.size = Pt(body_size)
        run.font.color.rgb = COLOR_BODY
        # East-asian font
        for r in p.runs:
            rPr = r._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
                rPr.append(ea)
        p.space_after = Pt(8)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def style_table_cell(cell, text, *, size=14, bold=False, color=COLOR_BODY, fill=None,
                     align=PP_ALIGN.LEFT, font=FONT_KR):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = ""
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
        rPr.append(ea)


def slide_table(prs, n, title, headers, rows, *, star=None, col_widths=None, body_size=14):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_top = Inches(1.55)
    table_h = Inches(min(5.0, 0.45 * n_rows + 0.2))
    tbl_shape = s.shapes.add_table(n_rows, n_cols, Inches(0.5), table_top, Inches(12.33), table_h)
    tbl = tbl_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)
    for ci, h in enumerate(headers):
        style_table_cell(tbl.cell(0, ci), h, size=body_size + 2, bold=True,
                         color=COLOR_TABLE_HEADER_FG, fill=COLOR_TABLE_HEADER_BG,
                         align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        fill = COLOR_TABLE_ROW_ALT if ri % 2 == 0 else None
        for ci, cell_text in enumerate(row):
            style_table_cell(tbl.cell(ri + 1, ci), str(cell_text), size=body_size, fill=fill)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_two_col(prs, n, title, left_title, left_body, right_title, right_body, *, star=None):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    # Two cards
    for col_i, (sub, body, color) in enumerate([
        (left_title, left_body, COLOR_ACCENT),
        (right_title, right_body, COLOR_STAR),
    ]):
        left = Inches(0.5 + col_i * 6.4)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.55), Inches(6.0), Inches(4.8))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        box.line.color.rgb = color
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)
        tf.margin_bottom = Inches(0.15)
        # Subtitle (header of card)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = sub
        run.font.name = FONT_KR
        run.font.size = Pt(22)
        run.font.bold = True
        run.font.color.rgb = color
        # Body lines
        for line in body.split('\n'):
            pp = tf.add_paragraph()
            pp.alignment = PP_ALIGN.LEFT
            r = pp.add_run()
            r.text = line
            r.font.name = FONT_KR
            r.font.size = Pt(16)
            r.font.color.rgb = COLOR_BODY
            pp.space_after = Pt(6)
            # East-asian
            rPr = r._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
                rPr.append(ea)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_code(prs, n, title, code, *, comment=None, star=None, code_size=14):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    # Code block
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.55), Inches(12.33), Inches(4.8))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CODE_BG
    box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    box.line.width = Pt(0.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    lines = code.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_CODE
        run.font.size = Pt(code_size)
        # color UE macros differently
        if any(m in line for m in ["UCLASS", "UPROPERTY", "UFUNCTION", "USTRUCT", "GENERATED_BODY"]):
            run.font.color.rgb = COLOR_ACCENT
            run.font.bold = True
        else:
            run.font.color.rgb = COLOR_TITLE
    if comment:
        add_textbox(s, Inches(0.5), Inches(6.0), Inches(12.33), Inches(0.6),
                    comment, size=14, italic=True, color=RGBColor(0x66, 0x66, 0x66),
                    align=PP_ALIGN.LEFT)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_experiment(prs, n, title, before, after, prediction):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    # Before / After labels
    add_textbox(s, Inches(0.5), Inches(1.55), Inches(6.0), Inches(0.4),
                "Before", size=18, bold=True, color=COLOR_ACCENT)
    add_textbox(s, Inches(6.83), Inches(1.55), Inches(6.0), Inches(0.4),
                "After (이걸 떼어보면…)", size=18, bold=True, color=COLOR_STAR)
    # Two code boxes
    for col_i, code in enumerate([before, after]):
        left = Inches(0.5 + col_i * 6.33)
        box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), Inches(6.0), Inches(2.5))
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_CODE_BG
        box.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.margin_top = Inches(0.1)
        tf.margin_bottom = Inches(0.1)
        for i, line in enumerate(code.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = line if line else " "
            run.font.name = FONT_CODE
            run.font.size = Pt(14)
            run.font.color.rgb = COLOR_TITLE
    # Prediction box
    pred = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.5), Inches(4.9), Inches(12.33), Inches(1.5))
    pred.fill.solid()
    pred.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xE1)
    pred.line.color.rgb = COLOR_STAR
    pred.line.width = Pt(1.5)
    tf = pred.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "🤔  예측 — "
    r.font.name = FONT_KR
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = COLOR_STAR
    r2 = p.add_run()
    r2.text = prediction
    r2.font.name = FONT_KR
    r2.font.size = Pt(20)
    r2.font.italic = True
    r2.font.color.rgb = COLOR_TITLE
    rPr = r2._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
        rPr.append(ea)
    add_page_number(s, n)
    return s


def slide_experiment_result(prs, n, title, observation, meaning, star):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    # Two observation/meaning boxes (stacked)
    boxes = [
        ("관찰 (실제 일어난 일)", observation, COLOR_ACCENT, RGBColor(0xE3, 0xF2, 0xFD)),
        ("의미 (왜 이렇게 되는가)", meaning, COLOR_STAR, RGBColor(0xFF, 0xEB, 0xEE)),
    ]
    for i, (label, body, line_color, fill_color) in enumerate(boxes):
        top = Inches(1.6 + i * 2.4)
        # Label
        add_textbox(s, Inches(0.5), top, Inches(3.5), Inches(0.45),
                    label, size=18, bold=True, color=line_color)
        # Body box
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(4.1), top, Inches(8.73), Inches(2.0))
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
        box.line.color.rgb = line_color
        box.line.width = Pt(1.5)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = body
        r.font.name = FONT_KR
        r.font.size = Pt(18)
        r.font.color.rgb = COLOR_TITLE
        rPr = r._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
            rPr.append(ea)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_flow(prs, n, title, flow_lines, star):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.5), Inches(1.7), Inches(10.33), Inches(4.5))
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_CODE_BG
    box.line.color.rgb = COLOR_ACCENT
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    for i, line in enumerate(flow_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.name = FONT_CODE
        r.font.size = Pt(18)
        r.font.color.rgb = COLOR_TITLE
        if line.strip() in ("↓", "▼"):
            r.font.color.rgb = COLOR_ACCENT
            r.font.bold = True
        p.space_after = Pt(4)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_callout(prs, n, title, body_lines, metaphor, star=None):
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    # Body bullets
    body_text = "\n".join(body_lines)
    add_textbox(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(3.5),
                body_text, size=22, color=COLOR_BODY, align=PP_ALIGN.LEFT)
    # Metaphor box
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(1.5), Inches(5.2), Inches(10.33), Inches(1.4))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE0)
    box.line.color.rgb = COLOR_STAR
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "💡  비유 — "
    r.font.name = FONT_KR
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = COLOR_STAR
    r2 = p.add_run()
    r2.text = metaphor
    r2.font.name = FONT_KR
    r2.font.size = Pt(20)
    r2.font.italic = True
    r2.font.color.rgb = COLOR_TITLE
    rPr = r2._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
        rPr.append(ea)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


def slide_four_box(prs, n, title, boxes, star=None):
    """boxes = [(label, body_text), ...] up to 4."""
    s = blank_slide(prs)
    add_title(s, title)
    add_underline(s)
    palette = [
        (COLOR_ACCENT, RGBColor(0xE3, 0xF2, 0xFD)),
        (COLOR_STAR, RGBColor(0xFF, 0xEB, 0xEE)),
        (COLOR_CHECK, RGBColor(0xE8, 0xF5, 0xE9)),
        (RGBColor(0x80, 0x57, 0xC2), RGBColor(0xF3, 0xE5, 0xF5)),
    ]
    box_w = Inches(6.0)
    box_h = Inches(2.3)
    positions = [(0.5, 1.6), (6.83, 1.6), (0.5, 4.1), (6.83, 4.1)]
    for i, (label, body) in enumerate(boxes[:4]):
        x, y = positions[i]
        line_color, fill_color = palette[i]
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = fill_color
        box.line.color.rgb = line_color
        box.line.width = Pt(2)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = label
        r.font.name = FONT_KR
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = line_color
        # body
        for line in body.split('\n'):
            pp = tf.add_paragraph()
            rr = pp.add_run()
            rr.text = line
            rr.font.name = FONT_KR
            rr.font.size = Pt(15)
            rr.font.color.rgb = COLOR_BODY
            pp.space_after = Pt(2)
            rPr = rr._r.get_or_add_rPr()
            ea = rPr.find(qn('a:ea'))
            if ea is None:
                ea = rPr.makeelement(qn('a:ea'), {'typeface': FONT_KR})
                rPr.append(ea)
    if star:
        add_star_footer(s, star)
    add_page_number(s, n)
    return s


# ===================== SLIDE CONTENT =====================

def build_slides(prs):
    n = 0

    # ----- Slide 1: Title -----
    n += 1
    slide_title_page(
        prs,
        title="언리얼 엔진 스터디 6주차",
        subtitle="매크로 한 글자가 결정하는 모든 것",
        footer="UCLASS  ·  UPROPERTY  ·  UFUNCTION  ·  BlueprintNativeEvent  ·  APickupItem",
    )

    # ----- Slide 2: 복습 1 — 6 매니저 -----
    n += 1
    slide_table(
        prs, n,
        title="복습 1 — 5주차의 6 매니저",
        headers=["매니저", "수명", "위치", "역할"],
        rows=[
            ["GameInstance", "앱 평생", "서버+클라", "레벨 전환에도 유지되는 전역 데이터"],
            ["GameMode", "레벨 단위", "서버 전용", "게임 룰, 승리/패배 판정"],
            ["GameState", "레벨 단위", "서버+모든 클라", "모두가 알아야 할 상태"],
            ["PlayerController", "접속 동안", "본인 클라+서버", "입력, Pawn 빙의, UI 소유"],
            ["PlayerState", "접속 동안", "서버+모든 클라", "HP, 점수, 보유 아이템"],
            ["HUD/UMG", "접속 동안", "본인 클라", "화면에 그리기"],
        ],
        col_widths=[2.5, 2.0, 3.0, 4.83],
        star="이 6개가 자리 잡혔으니, 오늘은 그 위에 *기능*을 얹는다",
    )

    # ----- Slide 3: 복습 2 — 룰 vs 상태 -----
    n += 1
    slide_table(
        prs, n,
        title="복습 2 — 룰 vs 상태",
        headers=["", "룰 (변하지 않음)", "상태 (실시간 변화)"],
        rows=[
            ["위치", "GameMode", "GameState / PlayerState"],
            ["예시", "MaxRespawnCount = 3", "RemainingRespawns = 2"],
            ["누가 바꾸나", "디자이너가 디테일 패널에서", "게임 진행이"],
        ],
        col_widths=[2.5, 4.5, 5.33],
        star="같은 '부활 횟수'라도 — 고정 한도는 GM, 현재 상황은 PS",
    )

    # ----- Slide 4: 복습 3 — 죽음 처리 흐름 -----
    n += 1
    slide_flow(
        prs, n,
        title="복습 3 — 죽음 처리는 PC가 책임진다",
        flow_lines=[
            "Timer (1s) → Character::OnTickDamage",
            "↓",
            "PS::ApplyDamage(10)  →  IsAlive() = false",
            "↓",
            "PC::HandlePlayerDeath()",
            "↓",
            "PS::DecrementRespawn()  →  UnPossess + Destroy  →  ScheduleRespawn",
            "↓",
            "0.5s 후  →  SpawnActor + Possess + PS::ResetHP()",
        ],
        star="Character는 '죽었다'만 알리고, 결정과 파괴는 PC가 한다 — 5주차 Q6의 답",
    )

    # ----- Slide 5: 이번 주의 질문 -----
    n += 1
    slide_quote(
        prs, n,
        title="이번 주의 질문 (한 줄로)",
        quote="C++과 블루프린트는 왜 둘 다 존재하는가?",
        subtext="한쪽이 무조건 좋다면 둘 다 있을 이유가 없다.\n그리고 — 매크로 한 글자(예: EditAnywhere)가 안 붙으면 정말로 디테일 패널에서 사라지는가?\n직접 떼어보면 답이 나온다.",
    )

    # ----- Slide 6: 오늘의 목표 -----
    n += 1
    slide_bullets(
        prs, n,
        title="오늘의 목표",
        bullets=[
            "✓ 언리얼식 C++ 매크로(UCLASS/UPROPERTY/UFUNCTION)가 왜 필요한지를 *리플렉션* 관점에서 설명할 수 있다",
            "✓ BlueprintCallable / BlueprintNativeEvent / BlueprintImplementableEvent 의 차이를 구분하고 쓸 수 있다",
            "✓ 매크로 옵션을 일부러 빠뜨려보면서 각 한 글자가 무엇을 결정하는지 손으로 체감한다",
            "✓ APickupItem 베이스를 상속해 BP 자식 클래스로 자기만의 아이템을 디자인한다",
            "★ (메타 목표) 5주차 매니저 구조의 가치를 *늦게* 깨닫는다",
        ],
        body_size=20,
    )

    # ----- Slide 7: 진행 방식 -----
    n += 1
    slide_table(
        prs, n,
        title="이번 주차 진행 방식",
        headers=["단계", "시간", "내용"],
        rows=[
            ["이론", "~30분", "매크로 4총사 + C++/BP 비교"],
            ["샘플 분석 + 매크로 실험", "~20분", "배포된 APickupItem 코드의 매크로를 직접 떼어본다"],
            ["라이브 데모", "~15분", "BP_SpeedBoost / BP_DamageBoost 시연"],
            ["아이디어 브레인스토밍", "~10분", "어떤 아이템을 만들고 싶은가"],
            ["직접 제작 + 발표", "~15분", "본인 아이디어 또는 가이드 옵션"],
        ],
        col_widths=[3.5, 1.8, 7.03],
        star="오늘은 듣는 시간보다 손으로 만지는 시간이 더 길다",
    )

    # ----- Slide 8: 넘어가기 전에 -----
    n += 1
    slide_bullets(
        prs, n,
        title="넘어가기 전에 — 우리가 매주 써온 그 매크로",
        bullets=[
            "4주차:  UPROPERTY(EditAnywhere) UInputAction* JumpAction;",
            "5주차:  UPROPERTY(EditAnywhere, BlueprintReadOnly, Category=\"Rules\") int32 MaxRespawnCount;",
            "★ 그런데 — 우리는 정말 이게 무엇인지 *알고* 썼나?",
            "오늘은 그 한 줄을 *직접 떼어보면서* 매크로의 정체를 파헤친다.",
        ],
        body_size=20,
    )

    # ----- Slide 9: 언리얼 C++의 정체 -----
    n += 1
    slide_callout(
        prs, n,
        title="언리얼 C++의 정체 — 표준 C++이 아니다",
        body_lines=[
            "•  컴파일러가 .cpp 만 보는 게 아니다",
            "•  UHT (Unreal Header Tool) 가 헤더의 매크로를 *먼저* 읽어 메타데이터(*.generated.h)를 만든다",
            "•  우리가 빌드하면 사실 두 단계: ① UHT  →  ② C++ 컴파일러",
        ],
        metaphor="매크로는 엔진에게 보내는 자기소개서. 안 쓰면 엔진은 이 변수가 누군지 모른다.",
    )

    # ----- Slide 10: UHT가 만들어주는 5가지 -----
    n += 1
    slide_bullets(
        prs, n,
        title="UHT 메타데이터로 가능해지는 것들",
        bullets=[
            "GC 추적 — 매크로 없는 포인터는 가비지 컬렉터가 모름 → 댕글링 가능",
            "직렬화 — .uasset 저장/로드, 세이브 게임",
            "에디터 노출 — 디테일 패널에서 만질 수 있게",
            "BP 연동 — BP 그래프에서 호출/접근 가능",
            "Replication — 멀티플레이 자동 복제",
        ],
        star="매크로 없는 멤버 = 엔진이 모르는 멤버",
        body_size=22,
    )

    # ----- Slide 11: 핵심 매크로 4총사 -----
    n += 1
    slide_table(
        prs, n,
        title="핵심 매크로 4총사",
        headers=["매크로", "용도", "자주 쓰는 옵션"],
        rows=[
            ["UCLASS()", "클래스 등록", "Blueprintable, BlueprintType, Abstract"],
            ["UPROPERTY()", "변수 등록 (GC 추적 포함)", "EditAnywhere, VisibleAnywhere, BlueprintReadWrite/ReadOnly, Replicated"],
            ["UFUNCTION()", "함수 등록", "BlueprintCallable, BlueprintNativeEvent, BlueprintImplementableEvent"],
            ["USTRUCT()", "구조체 등록", "BlueprintType"],
        ],
        col_widths=[2.5, 3.0, 6.83],
    )

    # ----- Slide 12: UPROPERTY 옵션 매트릭스 -----
    n += 1
    slide_table(
        prs, n,
        title="디자이너 vs 프로그래머 — UPROPERTY 옵션",
        headers=["옵션", "디테일 패널", "BP 그래프"],
        rows=[
            ["(옵션 없음)", "안 보임", "안 보임"],
            ["VisibleAnywhere", "보임 (회색, 수정 불가)", "안 보임"],
            ["EditAnywhere", "보임 (수정 가능)", "안 보임"],
            ["BlueprintReadOnly", "안 보임", "get만 가능"],
            ["BlueprintReadWrite", "안 보임", "get/set 가능"],
            ["EditAnywhere + BlueprintReadWrite", "둘 다 가능", "둘 다 가능"],
        ],
        col_widths=[5.0, 4.0, 3.33],
        star="옵션을 *조합*해서 노출 범위를 정밀하게 제어한다",
    )

    # ----- Slide 13: Native vs Implementable -----
    n += 1
    slide_two_col(
        prs, n,
        title="Native Event vs Implementable Event",
        left_title="BlueprintNativeEvent",
        left_body="C++에 *기본 동작 있음*\nBP에서 오버라이드 가능\n\n\"기본은 이렇게 작동하지만\n자식이 바꿀 수 있다\"\n\n예) APickupItem::OnPickedUp\n— 기본은 destroy + 효과 재생",
        right_title="BlueprintImplementableEvent",
        right_body="C++에 *구현 없음*\nBP가 반드시 채워야 한다\n\n\"이 부분은 디자이너 영역이다\"\n\n예) APickupItem::OnPlayPickupEffect\n— 시각/사운드 효과는 BP가 정함",
        star="비유 — 기본값이 채워진 메뉴 vs 빈칸의 설문지",
    )

    # ----- Slide 14: C++ vs Blueprint -----
    n += 1
    slide_table(
        prs, n,
        title="C++ vs Blueprint — 어디에 강한가",
        headers=["기준", "C++", "Blueprint"],
        rows=[
            ["실행 속도", "빠름", "느림 (VM)"],
            ["변경 주기", "컴파일 필요", "즉시"],
            ["디자이너 협업", "어려움", "쉬움"],
            ["복잡한 알고리즘", "강함", "노드 폭발"],
            ["데이터 튜닝", "약함", "강함 (에디터 즉시)"],
        ],
        col_widths=[3.5, 4.5, 4.33],
        star="한쪽이 무조건 좋으면 둘 다 있을 이유가 없다 — 강점이 정확히 반대",
    )

    # ----- Slide 15: 베이스 C++ + 변종 BP -----
    n += 1
    slide_callout(
        prs, n,
        title="현장 패턴 — 베이스는 C++, 변종은 BP",
        body_lines=[
            "•  Vampire Survivors 같은 게임에서 아이템이 30종이면? — 베이스 1개만 C++, 변종 30개는 모두 BP",
            "•  왜? — 데이터 튜닝이 잦은 변종을 BP로 빼면 컴파일 없이 즉시 밸런싱",
            "•  베이스의 핵심 알고리즘은 C++로 빠르게",
            "•  디자이너가 변종을 직접 만들고 만질 수 있다",
        ],
        metaphor="오늘 실습이 정확히 이 패턴 — 베이스 1 + 변종 3",
    )

    # ----- Slide 16: 디바이더 — APickupItem -----
    n += 1
    slide_divider(
        prs, n,
        title="APickupItem 코드 분석",
        subtitle="매크로를 일부러 떼어보자",
    )

    # ----- Slide 17: 배포 패키지 -----
    n += 1
    slide_table(
        prs, n,
        title="배포 패키지 — Week6_PickupSystem.zip",
        headers=["파일", "형태", "역할"],
        rows=[
            ["APickupItem.h/.cpp", "C++ 베이스 클래스", "모든 픽업의 부모 (Abstract)"],
            ["BP_HealthPotion", "BP 자식 (가장 단순)", "체력 즉시 회복 — 즉시 효과 패턴"],
            ["BP_DamageBoost", "BP 자식 (지속 효과)", "공격력 +50%, 10초 — Timer 패턴"],
            ["BP_SpeedBoost", "BP 자식 (시각 효과)", "이동속도 +50%, 5초 — OnPlayPickupEffect 활용"],
        ],
        col_widths=[3.5, 3.0, 5.83],
        star="세 샘플은 *즉시·지속·시각* 세 축의 패턴을 의도적으로 골랐다",
    )

    # ----- Slide 18: APickupItem 코드 -----
    n += 1
    slide_code(
        prs, n,
        title="APickupItem 베이스 클래스 — 전체 코드",
        code="""UCLASS(Abstract, Blueprintable)
class APickupItem : public AActor
{
    GENERATED_BODY()

public:
    APickupItem();

protected:
    UPROPERTY(VisibleAnywhere) USphereComponent* CollisionSphere;
    UPROPERTY(VisibleAnywhere) UStaticMeshComponent* MeshComp;

    UPROPERTY(EditAnywhere, BlueprintReadOnly, Category="Pickup")
    FText DisplayName;

    UFUNCTION(BlueprintNativeEvent, Category="Pickup")
    void OnPickedUp(AActor* Picker);

    UFUNCTION(BlueprintImplementableEvent, Category="Pickup")
    void OnPlayPickupEffect();

    virtual void BeginPlay() override;
};""",
        comment="이 코드 한 화면이 오늘 실습의 모든 출발점. 매크로 6줄에 주목.",
    )

    # ----- Slide 19: Abstract의 의미 -----
    n += 1
    slide_callout(
        prs, n,
        title="★ 왜 UCLASS(Abstract) 인가?",
        body_lines=[
            "•  APickupItem 자체로는 레벨에 못 놓는다 — 자식 클래스만 가능",
            "•  \"베이스 클래스는 인스턴스화하지 않는다\"는 의도를 매크로로 *강제*",
            "•  Blueprintable이 같이 있어야 BP 자식 생성이 가능 — 둘은 짝꿍",
        ],
        metaphor="엑셀 템플릿 — 그 자체로는 보고서가 아니다. 채워야 보고서가 된다.",
    )

    # ----- Slide 20: 디바이더 — 매크로 빠뜨려보기 -----
    n += 1
    slide_divider(
        prs, n,
        title="매크로 빠뜨려보기",
        subtitle="오늘의 진짜 핵심 — 손으로 한 번 느끼면 평생 안 잊는다",
    )

    # ----- Slide 21: 실험 1 셋업 -----
    n += 1
    slide_experiment(
        prs, n,
        title="실험 1 — UPROPERTY 자체를 떼면?",
        before='UPROPERTY(EditAnywhere, BlueprintReadOnly,\n          Category="Pickup")\nFText DisplayName;',
        after="FText DisplayName;",
        prediction="BP 자식을 컴파일하면 디테일 패널에서 어떻게 될까?",
    )

    # ----- Slide 22: 실험 1 결과 -----
    n += 1
    slide_experiment_result(
        prs, n,
        title="실험 1 결과",
        observation="BP 컴파일 후 디테일 패널에서 DisplayName 변수 자체가 *사라진다*. BP 그래프에서도 안 보인다.",
        meaning="GC 추적 X / 직렬화 X / 에디터 노출 X / BP 노출 X — 엔진에게 이 변수는 *존재하지 않는다*.",
        star="매크로 없는 멤버 = 엔진이 모르는 멤버",
    )

    # ----- Slide 23: 실험 2 셋업 -----
    n += 1
    slide_experiment(
        prs, n,
        title="실험 2 — EditAnywhere → VisibleAnywhere",
        before='UPROPERTY(EditAnywhere,\n          BlueprintReadOnly, ...)\nFText DisplayName;',
        after='UPROPERTY(VisibleAnywhere,\n          BlueprintReadOnly, ...)\nFText DisplayName;',
        prediction="디테일 패널에 보일까? 보이면 만질 수 있을까?",
    )

    # ----- Slide 24: 실험 2 결과 -----
    n += 1
    slide_experiment_result(
        prs, n,
        title="실험 2 결과",
        observation="디테일 패널에 보이긴 한다. 그런데 *회색* — 수정이 안 된다.",
        meaning="활용 케이스: \"보여주되 못 만지게\" — 디버그용 변수 노출 / 런타임 상태 모니터링.",
        star="Visible은 읽기 전용 노출, Edit은 쓰기 가능 노출",
    )

    # ----- Slide 25: 실험 3 셋업 -----
    n += 1
    slide_experiment(
        prs, n,
        title="실험 3 — BlueprintReadOnly 만 떼기",
        before='UPROPERTY(EditAnywhere,\n          BlueprintReadOnly,\n          Category="Pickup")\nFText DisplayName;',
        after='UPROPERTY(EditAnywhere,\n          Category="Pickup")\nFText DisplayName;',
        prediction="디테일 패널은 그대로일 텐데… BP 그래프에서 Get DisplayName 노드는?",
    )

    # ----- Slide 26: 실험 3 결과 -----
    n += 1
    slide_experiment_result(
        prs, n,
        title="실험 3 결과",
        observation="디테일 패널은 그대로 보인다. 그러나 BP 그래프에서 Get DisplayName 노드가 *나오지 않는다*.",
        meaning="디테일 패널 노출과 BP 그래프 노출은 *별개의 게이트*. EditAnywhere가 BP 노출을 자동으로 켜주지 않는다.",
        star="BP 노출은 추가 옵션(BlueprintReadOnly / ReadWrite)이 있어야 켜진다",
    )

    # ----- Slide 27: 실험 4 셋업 -----
    n += 1
    slide_experiment(
        prs, n,
        title="실험 4 — UFUNCTION() 자체를 떼면?",
        before='UFUNCTION(BlueprintNativeEvent,\n          Category="Pickup")\nvoid OnPickedUp(AActor* Picker);',
        after="void OnPickedUp(AActor* Picker);",
        prediction="BP 자식에서 OnPickedUp을 오버라이드하려고 하면?",
    )

    # ----- Slide 28: 실험 4 결과 -----
    n += 1
    slide_experiment_result(
        prs, n,
        title="실험 4 결과",
        observation="BP의 *Override 함수* 목록에 OnPickedUp 자체가 *나타나지 않는다*.",
        meaning="BP에서 호출/오버라이드되려면 UFUNCTION() 매크로가 *필수 게이트*. 함수도 변수와 동일한 원리.",
        star="변수는 UPROPERTY, 함수는 UFUNCTION — 게이트 두 종류",
    )

    # ----- Slide 29: 실험 5 셋업 -----
    n += 1
    slide_experiment(
        prs, n,
        title="실험 5 — Implementable 을 NativeEvent 로 바꾸면?",
        before="UFUNCTION(BlueprintImplementableEvent)\nvoid OnPlayPickupEffect();",
        after="UFUNCTION(BlueprintNativeEvent)\nvoid OnPlayPickupEffect();",
        prediction="컴파일은 통과할까?",
    )

    # ----- Slide 30: 실험 5 결과 -----
    n += 1
    slide_experiment_result(
        prs, n,
        title="실험 5 결과 — 컴파일 에러!",
        observation="컴파일이 *실패*한다. unresolved external symbol — OnPlayPickupEffect_Implementation() 의 *기본 구현*을 강제로 요구.",
        meaning="NativeEvent는 C++ 기본 구현이 *반드시 있어야* 한다 — 그게 정의. Implementable은 비워둘 수 있는 *유일한* 형태.",
        star="Implementable은 비워두는 게 명시적이다 — '여긴 디자이너 영역' 이라는 선언",
    )

    # ----- Slide 31: 실험 정리 -----
    n += 1
    slide_table(
        prs, n,
        title="매크로 실험 정리 — 한 줄이 결정한 것들",
        headers=["실험", "떼어본 것", "사라진 능력 (또는 추가 강제)"],
        rows=[
            ["1", "UPROPERTY 자체", "디테일·BP·GC·직렬화 모두"],
            ["2", "EditAnywhere → VisibleAnywhere", "수정 가능성"],
            ["3", "BlueprintReadOnly", "BP 그래프에서 read"],
            ["4", "UFUNCTION 자체", "BP에서 호출/오버라이드"],
            ["5", "Implementable → Native", "(반대로) 기본 구현 의무 *추가*"],
        ],
        col_widths=[1.0, 4.5, 6.83],
        star="매크로 한 줄 = 디테일 / BP / GC / 직렬화 / 컴파일 모두 결정",
    )

    # ----- Slide 32: 디바이더 — 라이브 데모 -----
    n += 1
    slide_divider(
        prs, n,
        title="라이브 데모 — BP 자식 만드는 표준 패턴",
        subtitle="강사가 처음부터 다시 만든다",
    )

    # ----- Slide 33: BP_SpeedBoost 의도 -----
    n += 1
    slide_bullets(
        prs, n,
        title="데모 ① BP_SpeedBoost — 시각 효과 패턴",
        bullets=[
            "★ 의도 — PIE에서 \"내가 빨라진다\"를 즉시 눈으로 본다 → 학습 동기",
            "다루는 패턴 — Native(OnPickedUp) + Implementable(OnPlayPickupEffect) 둘 다 사용",
            "결과물 — 5초간 이동속도 ×1.5, 그 후 자동 원복",
        ],
        body_size=22,
    )

    # ----- Slide 34: BP_SpeedBoost 만드는 순서 1/2 -----
    n += 1
    slide_bullets(
        prs, n,
        title="BP_SpeedBoost 만드는 순서 (1/2)",
        bullets=[
            "✓ Content Browser → 우클릭 → Blueprint Class → All Classes에서 APickupItem 선택",
            "✓ 이름 = BP_SpeedBoost",
            "✓ 디테일 패널: DisplayName = \"Speed Boost\", MeshComp에 임시 메쉬 (예: SM_ChamferCube)",
            "✓ BP 변수 추가: SpeedMultiplier (Float, 기본 1.5, EditAnywhere + BlueprintReadOnly)",
            "✓ BP 변수 추가: Duration (Float, 기본 5.0, EditAnywhere + BlueprintReadOnly)",
        ],
        body_size=20,
    )

    # ----- Slide 35: BP_SpeedBoost 그래프 -----
    n += 1
    slide_code(
        prs, n,
        title="BP_SpeedBoost 만드는 순서 (2/2)",
        code="""Event OnPickedUp (Picker) →
  Cast Picker To Week3Character →
    Get Character Movement →
      OriginalSpeed = MaxWalkSpeed
      Set MaxWalkSpeed ← OriginalSpeed × SpeedMultiplier →
    Set Timer By Event (Duration, OneShot) →
      [Custom Event] RestoreSpeed:
        Set MaxWalkSpeed ← OriginalSpeed
  Parent: OnPickedUp   ← 부모의 destroy + 효과 호출

OnPlayPickupEffect → Print String "Speed Up!" """,
        comment="Set Timer는 Picker(Character)의 World Timer Manager에 등록 — 픽업 액터가 destroy 되어도 타이머는 산다.",
        code_size=14,
    )

    # ----- Slide 36: BP_SpeedBoost 시연 포인트 -----
    n += 1
    slide_bullets(
        prs, n,
        title="BP_SpeedBoost 시연 포인트",
        bullets=[
            "✓ Native(OnPickedUp) + Implementable(OnPlayPickupEffect) 둘 다 사용",
            "✓ Timer 패턴 — 지속시간 있는 모든 아이템의 기본형",
            "✓ Picker(Character)의 World Timer Manager에 등록 → 픽업이 destroy 되어도 타이머는 산다",
        ],
        star="즉각적 피드백이 학습 동기를 만든다",
        body_size=22,
    )

    # ----- Slide 37: BP_DamageBoost 의도 -----
    n += 1
    slide_bullets(
        prs, n,
        title="데모 ② BP_DamageBoost — 매니저 연동 패턴",
        bullets=[
            "★ 의도 — 5주차의 PlayerState 구조와 *자연스럽게 합쳐지는* 걸 보여준다",
            "다루는 패턴 — Native만 사용 / Implementable은 *비워둠* (토론 유도)",
            "결과물 — 10초간 PS.DamageMultiplier = 1.5, 그 후 1.0으로 원복",
        ],
        body_size=22,
    )

    # ----- Slide 38: BP_DamageBoost 만드는 순서 -----
    n += 1
    slide_code(
        prs, n,
        title="BP_DamageBoost 만드는 순서",
        code="""1. APickupItem 상속 BP 생성 → BP_DamageBoost
2. 변수: BoostMultiplier (Float, 1.5),  Duration (Float, 10.0)
3. OnPickedUp 그래프:

     Picker → Get Player State → Cast Week3PlayerState (PS) →
       PS.DamageMultiplier ← BoostMultiplier →
     Set Timer By Event (Duration, OneShot) →
       [Custom Event] Reset:
         PS.DamageMultiplier ← 1.0
     Parent: OnPickedUp

4. OnPlayPickupEffect → 비워둠 (← 토론 포인트!)""",
        comment="Implementable Event를 비워두면 컴파일은 되지만 BP가 비어있다 — 의도적 빈 칸.",
        code_size=14,
    )

    # ----- Slide 39: 5주차 매니저 사후 정당화 -----
    n += 1
    slide_callout(
        prs, n,
        title="★ 5주차 분리의 가치가 여기서 보상받는다",
        body_lines=[
            "•  PS의 DamageMultiplier 한 줄만 바꿨는데 — *시스템이 알아서 반영*된다",
            "•  5주차 때 \"왜 이렇게 복잡하게 나누지?\" 라고 느꼈던 매니저 분리가 — 오늘 보상받음",
            "•  만약 모든 게 Character에 있었다면? — 같은 효과를 만들기 위해 Character 코드를 매번 수정해야 했을 것",
        ],
        metaphor="좋은 추상화는 *늦게* 보상받는다 — 5주차 분리는 6주차에서 답이 도착한다",
    )

    # ----- Slide 40: 두 데모 비교 -----
    n += 1
    slide_table(
        prs, n,
        title="두 데모 비교 — 같은 베이스, 다른 패턴",
        headers=["패턴", "SpeedBoost", "DamageBoost"],
        rows=[
            ["대상", "Character (이동 컴포넌트)", "PlayerState (데미지 배수)"],
            ["Implementable Event", "구현함 (Print String)", "비워둠"],
            ["5주차 매니저 활용", "X (Character 직접)", "O (PS 거쳐서)"],
            ["핵심 학습", "Timer 패턴", "매니저 연동 패턴"],
        ],
        col_widths=[3.5, 4.5, 4.33],
        star="같은 베이스에서 *완전히 다른 두 결*의 아이템이 나온다",
    )

    # ----- Slide 41: 아이디어 브레인스토밍 -----
    n += 1
    slide_four_box(
        prs, n,
        title="여러분이 만들고 싶은 아이템은? — 즉석 분류",
        boxes=[
            ("PS만 만지면 됨", "바로 만들 수 있음\n예: 무적, 데미지 부스트"),
            ("Timer 필요", "데모와 같은 패턴\n예: 일시 강화, 일시 무적"),
            ("9주차 이후 가능", "적이 있어야 작동\n예: 일정 범위 적 처치"),
            ("나중에", "시스템 추가 필요\n예: 인벤토리, 무기 슬롯"),
        ],
        star="내 아이디어가 어느 칸에 들어가는지가 — 5주차 구조의 *경계*를 보여준다",
    )

    # ----- Slide 42: 가이드 옵션 -----
    n += 1
    slide_bullets(
        prs, n,
        title="가이드 옵션 — 막힌 학생용",
        bullets=[
            "무적 아이템 (5초간 데미지 면역) → PS에 bIsInvincible 플래그 추가",
            "부활 횟수 +1 → PS의 RemainingRespawns ++",
            "체력 풀회복 + 속도 부스트 (복합 효과) → 두 데모를 합치기",
            "시한폭탄 회복 (3초 후 회복, 그 사이 도망쳐야) → Delay 패턴",
        ],
        body_size=20,
    )

    # ----- Slide 43: 제작 패턴 -----
    n += 1
    slide_bullets(
        prs, n,
        title="제작 패턴 — 5단계",
        bullets=[
            "✓ APickupItem 부모로 BP 생성",
            "✓ 필요한 변수 추가 (모두 EditAnywhere + BlueprintReadOnly)",
            "✓ OnPickedUp 오버라이드 — 효과 적용 + 필요시 Timer",
            "✓ OnPlayPickupEffect 구현 — 임시로 Print String",
            "✓ 레벨에 배치 → PIE 테스트",
        ],
        star="이 5단계가 모든 픽업 아이템의 *표준 레시피*",
        body_size=22,
    )

    # ----- Slide 44: 발표 -----
    n += 1
    slide_bullets(
        prs, n,
        title="발표 — 시간 남으면",
        bullets=[
            "2~3명이 자기가 만든 걸 시연",
            "질문: \"왜 이 아이템을 만들었나? 어디서 막혔나?\"",
            "다른 학생들의 아이디어를 듣는 게 *내 다음 아이템의 씨앗*",
        ],
        body_size=22,
    )

    # ----- Slide 45: 보너스 — 첫 게임 루프 -----
    n += 1
    slide_callout(
        prs, n,
        title="보너스 — 첫 게임 루프 완성",
        body_lines=[
            "•  5주차에서 만든 MaxRespawnCount = 1 설정",
            "•  본인이 만든 아이템 + 데모 아이템 + 데미지 시스템을 한 레벨에 모두 배치",
            "•  부활 + 데미지 + 회복 + 본인 제작 아이템이 *모두 작동하는 첫 순간*",
        ],
        metaphor="Vampire Survivors의 첫 게임 루프 완성 — 4-6주차 누적 결과",
    )

    # ----- Slide 46: 핵심 질문 -----
    n += 1
    slide_bullets(
        prs, n,
        title="핵심 질문 — 가져갈 것",
        bullets=[
            "1.  BlueprintNativeEvent 와 BlueprintImplementableEvent 를 언제 구분해서 쓰나?",
            "2.  디자이너가 포션 회복량을 자주 조절한다면 HealAmount 는 어떤 옵션 조합?",
            "3.  본인이 만든 아이템은 *PlayerState를 직접 만지나, Character를 만지나?* 어느 쪽이 더 깔끔할까?",
            "★ 4. 데모 아이템들은 PlayerState를 직접 만진다. 이게 정말 깔끔한 설계인가?",
        ],
        star="4번이 7주차로 가는 다리 — 답은 다음 주에",
        body_size=20,
    )

    # ----- Slide 47: 오늘 배운 것 -----
    n += 1
    slide_bullets(
        prs, n,
        title="오늘 배운 것",
        bullets=[
            "✓ 매크로 4총사 — UCLASS / UPROPERTY / UFUNCTION / USTRUCT",
            "✓ Native vs Implementable Event — 책임의 위치를 매크로로 표현",
            "✓ 베이스 C++ + 변종 BP — Vampire Survivors 패턴",
            "✓ 5주차 매니저 분리의 *사후 정당화* — DamageBoost가 PS만 만져도 작동",
            "✓ Vampire Survivors의 첫 완전한 게임 루프",
        ],
        star="한 줄 요약 — 매크로는 엔진에게 보내는 자기소개서",
        body_size=20,
    )

    # ----- Slide 48: 오늘 완성한 기능들 -----
    n += 1
    slide_bullets(
        prs, n,
        title="오늘 완성한 기능들 — 직접 돌려볼 수 있는 동작",
        bullets=[
            "✓ APickupItem 베이스 클래스 — Abstract + Blueprintable + 컴포넌트 + Native/Implementable Event",
            "✓ Week3PlayerState::Heal 함수 + DamageMultiplier 필드 추가",
            "✓ BP_HealthPotion — 즉시 회복",
            "✓ BP_SpeedBoost — 5초간 이동속도 +50%",
            "✓ BP_DamageBoost — 10초간 PS.DamageMultiplier = 1.5",
            "✓ 본인 제작 아이템 1종",
            "✓ 첫 완전한 게임 루프 (MaxRespawnCount=1 + 픽업 배치)",
        ],
        star="5주차에서 6 매니저가 자리를 잡았다면, 6주차는 그 위에 *기능*이 자라난 첫 주",
        body_size=18,
    )

    # ----- Slide 49: 다음 주 예고 -----
    n += 1
    slide_callout(
        prs, n,
        title="다음 주 예고 — 7주차",
        body_lines=[
            "•  마지막 핵심 질문의 답: \"PS를 직접 만지는 건 깔끔하지 않다\"",
            "•  HP·데미지 배수 등을 별도 *컴포넌트*로 추출 → 9주차의 적도 같은 컴포넌트를 쓴다",
            "•  HUD가 *폴링* 대신 *델리게이트로 구독*하는 구조로 전환",
        ],
        metaphor="오늘 만든 것을 다음 주에 한 번 더 추상화한다 — 좋은 추상화는 늦게 도착한다 (2)",
    )

    # ----- Slide 50: 정리 & Q&A -----
    n += 1
    s = blank_slide(prs)
    add_textbox(s, Inches(0.5), Inches(1.0), Inches(12.33), Inches(1.5),
                "정리 & Q&A", size=60, bold=True, color=COLOR_TITLE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.5), Inches(3.0), Inches(12.33), Inches(2.5),
                "다음 주 7주차: 5/21 (월) 18:30 — 컴포넌트 + 델리게이트\n"
                "본인 제작 아이템 코드를 GitHub PR로 공유 (선택)\n\n"
                "질문 받습니다.",
                size=24, color=COLOR_BODY, align=PP_ALIGN.CENTER)
    add_page_number(s, n)

    return n


def main():
    prs = Presentation(TEMPLATE)
    print(f"Loaded template: {len(list(prs.slides))} existing slides")
    remove_all_slides(prs)
    print(f"After strip: {len(list(prs.slides))} slides")
    total = build_slides(prs)
    print(f"Built {total} slides")
    prs.save(OUTPUT)
    print(f"Saved to: {OUTPUT}")


if __name__ == "__main__":
    main()
