#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_TEAL  = RGBColor(0x1B, 0x6B, 0x9A)
C_TEAL2 = RGBColor(0x0E, 0x45, 0x65)
C_GRAY  = RGBColor(0x55, 0x55, 0x55)
C_CYAN  = RGBColor(0x00, 0xE5, 0xFF)
FONT    = "맑은 고딕"

def I(v): return Inches(v)
def T(v): return Pt(v)

def new_prs():
    p = Presentation()
    p.slide_width  = I(13.33)
    p.slide_height = I(7.5)
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_text(sl, text, l, t, w, h, fs=22, bold=False,
             color=C_BLACK, align=PP_ALIGN.LEFT, italic=False):
    box = sl.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf  = box.text_frame
    tf.word_wrap = True
    pg  = tf.paragraphs[0]
    pg.alignment = align
    rn  = pg.add_run()
    rn.text = text
    rn.font.size  = T(fs)
    rn.font.bold  = bold
    rn.font.italic = italic
    rn.font.color.rgb = color
    rn.font.name  = FONT
    return box

def add_bullets(sl, items, l, t, w, h, fs=22, color=C_BLACK):
    """items: list of (indent_level 0|1, text, bold)"""
    box = sl.shapes.add_textbox(I(l), I(t), I(w), I(h))
    tf  = box.text_frame
    tf.word_wrap = True
    for i, (lvl, text, bold) in enumerate(items):
        pg = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        rn = pg.add_run()
        if text == "":
            rn.text = ""
        else:
            rn.text = ("• " if lvl == 0 else "    • ") + text
        rn.font.size  = T(fs if lvl == 0 else max(fs - 3, 16))
        rn.font.bold  = bold
        rn.font.color.rgb = color
        rn.font.name  = FONT
    return box

def add_rect(sl, l, t, w, h, fill=C_TEAL):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    sh = sl.shapes.add_shape(1, I(l), I(t), I(w), I(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh

def add_note(sl, note):
    try:
        sl.notes_slide.notes_text_frame.text = note
    except Exception:
        pass

# ---------------------------------------------------------------------------
# Slide type helpers
# ---------------------------------------------------------------------------

def mk_cover(prs, title, subtitle="", note=""):
    sl = blank(prs)
    add_text(sl, title, 0.5, 2.2, 12.33, 2.0,
             fs=52, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(sl, subtitle, 0.5, 4.3, 12.33, 1.0,
                 fs=24, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_note(sl, note)
    return sl

def mk_section(prs, title, subtitle="", note=""):
    sl = blank(prs)
    add_text(sl, title, 0.5, 2.4, 12.33, 2.0,
             fs=44, bold=True, color=C_BLACK, align=PP_ALIGN.CENTER)
    if subtitle:
        add_text(sl, subtitle, 0.5, 4.5, 12.33, 0.8,
                 fs=22, color=C_GRAY, align=PP_ALIGN.CENTER)
    add_note(sl, note)
    return sl

def mk_bullets(prs, title, items, title_fs=34, body_fs=21, note=""):
    sl = blank(prs)
    add_text(sl, title, 0.5, 0.2, 12.33, 1.0, fs=title_fs, bold=True)
    add_rect(sl, 0.5, 1.1, 12.33, 0.05, fill=C_TEAL)
    add_bullets(sl, items, 0.5, 1.3, 12.3, 5.8, fs=body_fs)
    add_note(sl, note)
    return sl

def mk_teal(prs, title, items, note=""):
    sl = blank(prs)
    add_rect(sl, 0.15, 0.15, 13.03, 7.2, fill=C_TEAL)
    add_rect(sl, 0.35, 0.35, 12.63, 6.8, fill=C_TEAL2)
    add_text(sl, title, 0.55, 0.4, 12.2, 1.0,
             fs=30, bold=True, color=C_WHITE)
    add_bullets(sl, items, 0.9, 1.5, 11.8, 5.5, fs=20, color=C_WHITE)
    add_note(sl, note)
    return sl

def mk_two_col(prs, title, left_title, left_items,
               right_title, right_items, body_fs=20, note=""):
    sl = blank(prs)
    add_text(sl, title, 0.5, 0.2, 12.33, 1.0, fs=34, bold=True)
    add_rect(sl, 0.5, 1.1, 12.33, 0.05, fill=C_TEAL)
    if left_title:
        add_rect(sl, 0.5, 1.25, 5.8, 0.5, fill=C_TEAL)
        add_text(sl, left_title, 0.5, 1.25, 5.8, 0.5,
                 fs=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_bullets(sl, left_items, 0.5, 1.85, 5.8, 5.3, fs=body_fs)
    if right_title:
        add_rect(sl, 6.83, 1.25, 5.8, 0.5, fill=C_TEAL)
        add_text(sl, right_title, 6.83, 1.25, 5.8, 0.5,
                 fs=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_bullets(sl, right_items, 6.83, 1.85, 5.8, 5.3, fs=body_fs)
    add_note(sl, note)
    return sl

def mk_step(prs, section_title, step_num, step_title, items, tip="", note=""):
    sl = blank(prs)
    add_text(sl, section_title, 0.5, 0.15, 12.3, 0.7, fs=24, bold=True, color=C_GRAY)
    add_rect(sl, 0.5, 0.95, 1.5, 0.65, fill=C_TEAL)
    add_text(sl, f"Step {step_num}", 0.5, 0.95, 1.5, 0.65,
             fs=19, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, step_title, 2.1, 0.95, 10.7, 0.65, fs=22, bold=True)
    add_rect(sl, 0.5, 1.65, 12.3, 0.04, fill=C_TEAL)
    add_bullets(sl, items, 0.5, 1.8, 12.3, 4.6, fs=20)
    if tip:
        add_rect(sl, 0.5, 6.6, 12.3, 0.7, fill=C_TEAL2)
        add_text(sl, "★ " + tip, 0.65, 6.6, 12.0, 0.7,
                 fs=16, color=C_WHITE, italic=True)
    add_note(sl, note)
    return sl

# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------

def build(out_path):
    prs = new_prs()

    # ── Slide 1: Cover ──────────────────────────────────────────────────────
    mk_cover(prs,
             "언리얼 엔진 스터디 4주차",
             "PlayerController · PlayerState · Enhanced Input · UMG 입문",
             note="오늘의 큰 그림: Controller가 Pawn을 '조종'하는 구조를 이해하고,\n"
                  "스프린트·HP 바·피해 처리까지 VS-라이크 프로젝트를 함께 만들어 봅니다.")

    # ── Slide 2: 오늘의 목표 ────────────────────────────────────────────────
    mk_bullets(prs, "오늘의 목표",
               [(0, "PlayerController / PlayerState의 역할과 차이를 설명할 수 있다", False),
                (0, "Enhanced Input에서 InputMappingContext · InputAction을 연결할 수 있다", False),
                (0, "스프린트 기능을 C++로 구현하고 블루프린트에서 확인할 수 있다", False),
                (0, "UMG로 간단한 HP 바 위젯을 만들고 화면에 표시할 수 있다", False),
                (0, "PlayerState에 HP를 저장하고 데미지 로직을 연결할 수 있다", False)],
               note="목표를 함께 읽으면서 오늘 무엇을 만들지 그림을 그려줍니다.\n"
                    "VS-라이크 프로젝트 맥락: 매주 한 기능씩 쌓아가는 누적 프로젝트입니다.")

    # ── Slide 3: VS-라이크 프로젝트 소개 ────────────────────────────────────
    mk_bullets(prs, "우리가 만드는 것 — Vampire Survivors-like",
               [(0, "탑다운 시점 / 단일 플레이어 / 웨이브 생존", False),
                (0, "이번 주 추가: 스프린트 · HP 바 UI · 데미지 처리", False),
                (1, "다음 주: 적 스폰 시스템 + 피격 연출 (5주차)", False),
                (1, "이후: 경험치·레벨업·아이템 드롭 (6–8주차)", False),
                (0, "핵심 원칙: 매주 실행 가능한 빌드 유지", False)],
               note="프로젝트 전체 로드맵을 한 번 보여주면 학생들이 '왜 이걸 배우나'를 이해합니다.\n"
                    "실제 게임처럼 느껴지도록 이름/아이콘을 붙여주세요.")

    # ── Slide 4: 오늘 시간 배분 ─────────────────────────────────────────────
    mk_bullets(prs, "오늘 시간 배분 (120분)",
               [(0, "에디터 기초 숙달 (20분)", False),
                (1, "뷰포트 조작 / 기즈모 / 플레이 모드", False),
                (0, "이론 A — PlayerController (10분)", False),
                (0, "이론 B — PlayerState (8분)", False),
                (0, "이론 C — Enhanced Input 심화 (12분)", False),
                (0, "실습 1: 스프린트 구현 (15분)", False),
                (0, "실습 2: HP 바 위젯 (20분)", False),
                (0, "실습 3: PlayerState HP + 데미지 (20분)", False),
                (0, "보너스 / Q&A (5분+)", False)],
               note="시간 배분은 유연하게 운영합니다. 실습이 밀리면 보너스를 생략합니다.")

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION: 에디터 기초
    # ═══════════════════════════════════════════════════════════════════════
    mk_section(prs, "에디터 기초 숙달", "뷰포트 · 기즈모 · 블루프린트 에디터",
               note="실습 전에 에디터 조작에 익숙해지는 시간입니다. 손이 느리면 실습이 막힙니다.")

    # ── Slide 6: 뷰포트 조작 ────────────────────────────────────────────────
    mk_bullets(prs, "뷰포트 조작 — 마우스 단축키",
               [(0, "우클릭 드래그 → 카메라 회전 (FPS 시점)", False),
                (0, "우클릭 + W/A/S/D → 1인칭 이동", False),
                (0, "Alt + 좌클릭 드래그 → 오브젝트 주위 공전", False),
                (0, "Alt + 우클릭 드래그 → 줌인/아웃", False),
                (0, "F → 선택 오브젝트 포커스", False),
                (0, "End → 바닥에 붙이기 (Snap to Floor)", False),
                (1, "배치 후 필수! 바닥 아래로 떨어지는 버그 방지", False)],
               note="직접 조작해 보게 하세요. 마우스 감각이 없으면 이후 실습이 불편합니다.")

    # ── Slide 7: 기즈모와 스냅 ─────────────────────────────────────────────
    mk_bullets(prs, "기즈모와 스냅",
               [(0, "W / E / R → 이동 · 회전 · 스케일 기즈모 전환", False),
                (0, "Space → 로컬/월드 좌표계 전환", False),
                (0, "그리드 스냅: 뷰포트 우상단 자석 아이콘", False),
                (1, "이동 스냅: 10 cm 단위 권장", False),
                (1, "회전 스냅: 5° 또는 15° 권장", False),
                (0, "Ctrl+Z / Ctrl+Y → 실행 취소 / 재실행 (에디터 내)", False)],
               note="기즈모 색상: 빨강=X, 초록=Y, 파랑=Z — UE5 표준")

    # ── Slide 8: 플레이 모드 ────────────────────────────────────────────────
    mk_bullets(prs, "플레이 모드 3가지",
               [(0, "Play (뷰포트 안) — Alt+P", False),
                (1, "현재 뷰포트에서 즉시 플레이. 가장 빠름", False),
                (0, "Play (새 창) — 별도 창에서 실행", False),
                (1, "해상도 테스트 / 멀티플레이 시뮬레이션에 유용", False),
                (0, "Simulate — 플레이어 없이 월드만 시뮬레이션", False),
                (1, "AI·물리 동작 확인, 플레이어 컨트롤 없음", False),
                (0, "Eject (F8) — 플레이 중 에디터 조작 활성화", False)],
               note="Eject는 런타임 오브젝트 값을 실시간으로 보기 좋습니다.")

    # ── Slide 9: 블루프린트 에디터 레이아웃 ──────────────────────────────────
    mk_bullets(prs, "블루프린트 에디터 레이아웃",
               [(0, "Viewport — 컴포넌트 배치·시각 확인", False),
                (0, "Construction Script — Begin Play 이전 초기화 로직", False),
                (0, "Event Graph — 게임플레이 이벤트 로직", False),
                (0, "My Blueprint 패널 — 변수 · 함수 · 매크로 목록", False),
                (0, "Details 패널 — 선택 노드/컴포넌트 속성", False),
                (0, "우클릭 → 노드 검색 팝업 (Context Sensitive 체크 해제 시 전체)", False)],
               note="'Context Sensitive'를 끄면 모든 노드가 보입니다. 막히면 먼저 꺼보세요.")

    # ── Slide 10: 프로젝트 탐색 ─────────────────────────────────────────────
    mk_bullets(prs, "프로젝트 탐색 — Content Browser & Outliner",
               [(0, "Content Browser — 모든 에셋 관리", False),
                (1, "필터 버튼으로 Blueprint / Texture / Sound 등 분류", False),
                (1, "우클릭 → Asset Actions → Rename (참조 자동 갱신)", False),
                (0, "World Outliner — 현재 레벨 오브젝트 트리", False),
                (1, "부모-자식 관계, 레이어, 검색 지원", False),
                (0, "Output Log — 런타임 로그 & 에러 확인", False),
                (1, "UE_LOG / PrintString 출력 확인 필수", False)],
               note="Rename을 Content Browser에서 해야 참조가 깨지지 않습니다. 파일 탐색기로 하면 안 됩니다!")

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION: 이론 A — PlayerController
    # ═══════════════════════════════════════════════════════════════════════
    mk_section(prs, "이론 A — PlayerController",
               "입력을 받아 폰을 조종하는 '사람의 의지'",
               note="핵심 비유: Controller = 사람의 두뇌, Pawn = 몸. 두뇌는 몸 없이도 존재합니다.")

    # ── Slide 12: 왜 PlayerController가 필요한가? ─────────────────────────
    mk_bullets(prs, "왜 PlayerController가 필요한가?",
               [(0, "3주차까지: 캐릭터 내부에 입력 처리 코드", False),
                (1, "캐릭터가 죽으면 입력도 함께 사라짐", False),
                (1, "탑다운 · RTS처럼 폰을 교체하면? 입력 코드도 이사해야?", False),
                (0, "해결: 입력·카메라·UI를 PlayerController로 분리", False),
                (1, "폰이 죽어도 Controller는 살아있음 → Respawn 구현 가능", False),
                (1, "Possess() 한 번으로 다른 폰 조종 가능", False)],
               note="학생들에게 '캐릭터가 죽을 때 어떻게 되나?' 질문으로 문제 인식을 먼저 끌어내세요.")

    # ── Slide 13: PlayerController 개념 ─────────────────────────────────────
    mk_bullets(prs, "PlayerController — 핵심 개념",
               [(0, "APlayerController extends AController", False),
                (0, "한 플레이어당 하나, 레벨 전환에도 유지됨", False),
                (0, "Possess(APawn*) — 폰에 빙의", False),
                (0, "UnPossess() — 폰 해제", False),
                (0, "GetPawn() — 현재 빙의 중인 폰 반환", False),
                (0, "입력 처리: SetupInputComponent 또는 Enhanced Input 바인딩", False),
                (0, "UI 관리: CreateWidget + AddToViewport 권장 위치", False)],
               note="Possess 시 OnPossess / UnPossess 이벤트가 호출됩니다. 이후 실습에서 확인합니다.")

    # ── Slide 14: Possess 흐름 ───────────────────────────────────────────────
    mk_bullets(prs, "Possess 실행 흐름",
               [(0, "GameMode::SpawnDefaultPawnFor() → Pawn 스폰", False),
                (0, "GameMode::RestartPlayer() → Possess 호출", False),
                (0, "APlayerController::Possess(pawn)", False),
                (1, "Pawn::PossessedBy(controller) 호출", False),
                (1, "OnPossess 이벤트 발생 (BP에서 오버라이드 가능)", False),
                (0, "이후 PlayerController가 입력을 해당 Pawn으로 전달", False)],
               note="다이어그램 그리며 설명하면 효과적입니다. 화이트보드 활용 권장.")

    # ── Slide 15: One Controller, One Character ─────────────────────────────
    mk_teal(prs, "핵심 규칙 — One Controller, One Pawn",
            [(0, "PlayerController는 한 번에 하나의 Pawn만 Possess 할 수 있다", True),
             (0, "멀티플레이: 각 플레이어마다 별도 PlayerController 인스턴스", False),
             (0, "AI: AIController가 같은 역할 수행 (PlayerController 사용 안 함)", False),
             (0, "Possess 교체 시 이전 Pawn은 자동 UnPossess", False),
             (0, "▶ 오늘 실습에서 BP_ThirdPersonCharacter를 그대로 사용", False)],
            note="AI와 Player의 Controller 타입이 다르다는 점을 강조하세요.")

    # ── Slide 16: 데모 — GetPlayerController ────────────────────────────────
    mk_bullets(prs, "데모: PlayerController 확인하기",
               [(0, "레벨 블루프린트 → Event BeginPlay", False),
                (0, "Get Player Controller (Index 0) → Print String", False),
                (1, "출력: BP_ThirdPersonPlayerController_C_0 형태", False),
                (0, "Get Controlled Pawn → Print String", False),
                (1, "출력: BP_ThirdPersonCharacter_C_0 형태", False),
                (0, "플레이 후 Output Log 에서도 확인 가능", False)],
               note="직접 해보게 하세요. 5분 이내 확인 가능한 간단한 확인 실습입니다.")

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION: 이론 B — PlayerState
    # ═══════════════════════════════════════════════════════════════════════
    mk_section(prs, "이론 B — PlayerState",
               "플레이어 데이터의 중앙 저장소",
               note="Controller = 의지, Pawn = 몸, PlayerState = 명찰(이름·HP·점수)")

    # ── Slide 18: PlayerState 개념 ──────────────────────────────────────────
    mk_bullets(prs, "PlayerState — 개념",
               [(0, "APlayerState — 플레이어 고유 데이터 저장", False),
                (1, "HP, 점수, 이름, 레벨, 경험치 등", False),
                (0, "PlayerController · Pawn 모두에서 접근 가능", False),
                (1, "GetPlayerState<AMyPlayerState>()", False),
                (0, "Pawn이 파괴되어도 PlayerState는 살아있음", False),
                (0, "멀티플레이: 서버에서 복제(Replicate)되어 모든 클라이언트에 동기화", False),
                (0, "오늘 목표: BP_PlayerState에 HP 변수 추가 → 데미지 처리", False)],
               note="멀티플레이 복제 개념은 간단히만 언급, 오늘 핵심은 단일 플레이어 구현입니다.")

    # ── Slide 19: GameState vs PlayerState ──────────────────────────────────
    mk_two_col(prs, "GameState vs PlayerState",
               "GameState", [
                   (0, "게임 전체 데이터", False),
                   (0, "웨이브 번호, 남은 시간", False),
                   (0, "모든 플레이어 공유", False),
                   (0, "GameMode가 관리", False),
                   (0, "", False),
                   (0, "AGameStateBase 상속", False),
               ],
               "PlayerState", [
                   (0, "개별 플레이어 데이터", False),
                   (0, "HP, 점수, 킬 수", False),
                   (0, "해당 플레이어 전용", False),
                   (0, "PlayerController가 소유", False),
                   (0, "", False),
                   (0, "APlayerState 상속", False),
               ],
               note="양쪽 다 WorldSettings에서 클래스 지정합니다.")

    # ── Slide 20: VS-라이크에서의 PlayerState ────────────────────────────────
    mk_teal(prs, "VS-라이크 프로젝트 — PlayerState 설계",
            [(0, "BP_PlayerState 변수", True),
             (1, "MaxHP (Float) = 100.0", False),
             (1, "CurrentHP (Float) = 100.0", False),
             (1, "Score (Integer) = 0", False),
             (0, "함수 (오늘 실습 3에서 구현)", True),
             (1, "TakeDamage(amount) — CurrentHP 감소 + 0 클램프", False),
             (1, "IsAlive() — CurrentHP > 0 반환", False),
             (0, "향후 추가: 경험치 / 레벨 / 아이템 보유 목록", False)],
            note="설계를 미리 보여주면 학생들이 실습 목적을 명확히 이해합니다.")

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION: 이론 C — Enhanced Input
    # ═══════════════════════════════════════════════════════════════════════
    mk_section(prs, "이론 C — Enhanced Input 심화",
               "InputMappingContext · InputAction · 런타임 교체",
               note="3주차에 기초를 배웠습니다. 오늘은 새 액션 추가 + 런타임 IMC 교체를 다룹니다.")

    # ── Slide 22: 레거시 Input의 문제 ───────────────────────────────────────
    mk_bullets(prs, "레거시 Input 시스템의 한계",
               [(0, "Project Settings → Input에 Axis/Action 하드코딩", False),
                (1, "모든 캐릭터가 같은 입력 설정을 공유", False),
                (0, "런타임 키 리맵핑 불가", False),
                (0, "입력 우선순위 제어 어려움", False),
                (0, "Modifier(감도 조정 등) 표준 방법 없음", False)],
               note="레거시를 직접 보여주면서 '이게 왜 불편한지' 느끼게 해주세요.")

    # ── Slide 23: Enhanced Input 구조 ───────────────────────────────────────
    mk_bullets(prs, "Enhanced Input — 3가지 핵심 에셋",
               [(0, "InputAction (IA_*) — '무엇'을 하는지", False),
                (1, "IA_Move, IA_Jump, IA_Sprint, IA_DebugDamage 등", False),
                (0, "InputMappingContext (IMC_*) — 키와 액션 매핑", False),
                (1, "IMC_Default: WASD → IA_Move, Space → IA_Jump", False),
                (1, "IMC_Vehicle: 별도 키 레이아웃", False),
                (0, "InputModifier — 값 변환 (Negate, Swizzle, DeadZone 등)", False),
                (1, "Swizzle YXZ: Y축 입력을 X로 → 전후 이동에 사용", False)],
               note="IA는 '행위', IMC는 '키 번역표' — 이 비유를 반복해서 사용하세요.")

    # ── Slide 24: AddMappingContext ──────────────────────────────────────────
    mk_bullets(prs, "런타임 IMC 추가/제거",
               [(0, "BeginPlay 또는 Possess 시점에 IMC 등록", False),
                (0, "C++ 코드:", False),
                (1, "UEnhancedInputLocalPlayerSubsystem* Sub =", False),
                (1, "  ULocalPlayer::GetSubsystem<...>(PlayerController);", False),
                (1, "Sub->AddMappingContext(IMC_Default, Priority=0);", False),
                (0, "Priority가 높은 IMC가 같은 키를 먼저 처리", False),
                (1, "예: 탈것 탑승 시 IMC_Vehicle(Priority=1) 추가", False),
                (0, "RemoveMappingContext(IMC)로 해제 가능", False)],
               note="Priority 개념을 UI 레이어에 비유하면 이해가 쉽습니다.")

    # ── Slide 25: BindAction 패턴 ───────────────────────────────────────────
    mk_bullets(prs, "IA 바인딩 패턴 (C++)",
               [(0, "SetupPlayerInputComponent에서:", False),
                (1, "UEnhancedInputComponent* EIC = Cast<UEnhancedInputComponent>(InputComponent);", False),
                (0, "Started 트리거 — 키를 처음 누르는 순간 1회", False),
                (1, "EIC->BindAction(IA_Jump, ETriggerEvent::Started, this, &AMyChar::Jump);", False),
                (0, "Triggered 트리거 — 키를 누르는 동안 매 프레임", False),
                (1, "EIC->BindAction(IA_Move, ETriggerEvent::Triggered, this, &AMyChar::Move);", False),
                (0, "Completed 트리거 — 키를 뗄 때 1회", False),
                (1, "EIC->BindAction(IA_Sprint, ETriggerEvent::Completed, this, &AMyChar::StopSprint);", False)],
               note="Started/Triggered/Completed 구분을 헷갈리면 스프린트 버그가 납니다. 강조하세요.")

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION: 실습 1 — 스프린트
    # ═══════════════════════════════════════════════════════════════════════
    mk_section(prs, "실습 1 — 스프린트 구현",
               "IA_Sprint + MaxWalkSpeed 변경 (15분)",
               note="가장 단순한 실습. Enhanced Input 흐름을 처음 손으로 만들어 보는 경험입니다.")

    # ── Steps ────────────────────────────────────────────────────────────────
    mk_step(prs, "실습 1: 스프린트",
            1, "IA_Sprint InputAction 에셋 생성",
            [(0, "Content Browser → 우클릭 → Input → Input Action", False),
             (0, "이름: IA_Sprint", False),
             (0, "Value Type: Digital (bool)", False),
             (0, "저장 (Ctrl+S)", False)],
            tip="Value Type을 'Digital'로 설정하지 않으면 Triggered가 매 프레임 호출됩니다.",
            note="에셋 이름 규칙: IA_ 접두사. Content/Input/ 폴더에 정리하면 깔끔합니다.")

    mk_step(prs, "실습 1: 스프린트",
            2, "IMC_Default에 Shift 키 매핑",
            [(0, "IMC_Default 열기 → Mappings + 버튼", False),
             (0, "Action: IA_Sprint 선택", False),
             (0, "키: Left Shift", False),
             (0, "저장", False)],
            note="IMC에서 키를 바꾸기만 하면 C++ 코드 변경 없이 키 리맵핑이 됩니다.")

    mk_step(prs, "실습 1: 스프린트",
            3, "Week3Character.h 에 함수·변수 추가",
            [(0, "UPROPERTY(EditAnywhere) float SprintSpeed = 800.f;", False),
             (0, "UPROPERTY(EditAnywhere) float WalkSpeed   = 500.f;", False),
             (1, "(기본값은 CharacterMovement의 MaxWalkSpeed와 맞추세요)", False),
             (0, "void Sprint(const FInputActionValue& Value);", False),
             (0, "void StopSprint(const FInputActionValue& Value);", False)],
            note="UPROPERTY로 노출하면 블루프린트에서 기본값을 편집할 수 있습니다.")

    mk_step(prs, "실습 1: 스프린트",
            4, "Week3Character.cpp 구현 + 바인딩",
            [(0, "Sprint(): GetCharacterMovement()->MaxWalkSpeed = SprintSpeed;", False),
             (0, "StopSprint(): GetCharacterMovement()->MaxWalkSpeed = WalkSpeed;", False),
             (0, "SetupPlayerInputComponent에 바인딩 추가:", False),
             (1, "EIC->BindAction(SprintAction, ETriggerEvent::Started,   this, &Sprint);", False),
             (1, "EIC->BindAction(SprintAction, ETriggerEvent::Completed,  this, &StopSprint);", False),
             (0, "SprintAction 변수(UPROPERTY)에 IA_Sprint 에셋 연결 (BP에서)", False)],
            tip="빌드 후 BP_ThirdPersonCharacter 디테일 패널에서 SprintAction 슬롯이 보여야 합니다.",
            note="컴파일 에러 시: Value Type Digital이면 FInputActionValue를 bool로 캐스팅할 수 있습니다.")

    # ── 실습 1 결과 확인 ─────────────────────────────────────────────────────
    mk_teal(prs, "실습 1 체크포인트",
            [(0, "✅ Shift 누르는 동안 이동 속도가 빨라진다", False),
             (0, "✅ Shift를 떼면 원래 속도로 돌아온다", False),
             (0, "✅ BP_ThirdPersonCharacter 디테일에서 SprintSpeed 값 변경 가능", False),
             (0, "✅ Output Log에 에러 없음", False),
             (0, "막히면? → 강사 호출 또는 옆 자리와 비교", False)],
            note="체크포인트에서 손을 들어 확인받게 하세요. 막힌 학생 파악 후 넘어갑니다.")

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION: 실습 2 — HP 바 위젯
    # ═══════════════════════════════════════════════════════════════════════
    mk_section(prs, "실습 2 — HP 바 UMG 위젯",
               "WBP_PlayerHUD 생성 + 화면에 표시 (20분)",
               note="UMG 첫 경험입니다. UI 레이어 개념과 ProgressBar 바인딩을 배웁니다.")

    mk_step(prs, "실습 2: HP 바 위젯",
            1, "WBP_PlayerHUD 위젯 블루프린트 생성",
            [(0, "Content Browser → 우클릭 → User Interface → Widget Blueprint", False),
             (0, "이름: WBP_PlayerHUD", False),
             (0, "기본 클래스: UserWidget", False),
             (0, "더블 클릭 → Widget Designer 열기", False)],
            note="Widget Blueprint는 UI 패널입니다. 이후 C++ Widget 클래스로 전환할 수 있습니다.")

    mk_step(prs, "실습 2: HP 바 위젯",
            2, "UI 레이아웃 구성",
            [(0, "Palette → Canvas Panel 추가 (루트)", False),
             (0, "좌하단에 Vertical Box 배치", False),
             (0, "Vertical Box 안에 Text Block 추가: 'HP'", False),
             (0, "Text Block 아래에 Progress Bar 추가", False),
             (1, "Size: 200×20, Fill Color: 빨간색(R=1,G=0,B=0)", False),
             (0, "Is Variable 체크 → 변수 이름: HPBar", False)],
            tip="'Is Variable'을 체크해야 그래프에서 HPBar 변수에 접근할 수 있습니다.",
            note="앵커(Anchor)를 좌하단으로 설정하면 해상도 변경 시 위치가 유지됩니다.")

    mk_step(prs, "실습 2: HP 바 위젯",
            3, "Percent 바인딩 + PlayerController에서 위젯 생성",
            [(0, "HPBar → Percent → Bind → Create Binding", False),
             (0, "바인딩 함수 안: Get Player State → Cast to BP_PlayerState", False),
             (1, "→ CurrentHP / MaxHP → Return (Float)", False),
             (0, "BP_ThirdPersonPlayerController → Begin Play:", False),
             (1, "Create Widget (WBP_PlayerHUD) → Add to Viewport", False),
             (1, "위젯 레퍼런스 변수(HUDWidget)에 저장", False)],
            tip="Cast 실패 시 분기를 잊으면 게임이 크래시됩니다. IsValid로 항상 확인!",
            note="Create Widget의 Owning Player를 Get Player Controller(0)으로 설정하세요.")

    # ── 실습 2 결과 확인 ─────────────────────────────────────────────────────
    mk_teal(prs, "실습 2 체크포인트",
            [(0, "✅ 게임 시작 시 좌하단에 빨간 프로그레스 바가 보인다", False),
             (0, "✅ 바 Percent가 1.0 (100%)으로 꽉 차 있다", False),
             (0, "✅ 위젯이 PIE 종료 후 다시 시작해도 표시된다", False),
             (0, "아직 HP가 고정값 — 실습 3에서 데미지를 연결합니다", False)],
            note="바가 안 보이면: Add to Viewport가 호출됐는지, ZOrder가 0인지 확인하세요.")

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION: 실습 3 — PlayerState HP + 데미지
    # ═══════════════════════════════════════════════════════════════════════
    mk_section(prs, "실습 3 — PlayerState HP + 데미지",
               "BP_PlayerState · TakeDamage · IA_DebugDamage (20분)",
               note="세 실습 중 가장 복잡합니다. PlayerState 연동과 데미지 흐름을 이해하는 게 목표입니다.")

    mk_step(prs, "실습 3: PlayerState HP + 데미지",
            1, "BP_PlayerState 생성 및 변수 추가",
            [(0, "Content Browser → Blueprint Class → PlayerState 검색 → 선택", False),
             (0, "이름: BP_PlayerState", False),
             (0, "변수 추가: MaxHP (Float, 100.0), CurrentHP (Float, 100.0)", False),
             (0, "Event BeginPlay: CurrentHP = MaxHP", False),
             (0, "World Settings → Game Mode → Player State Class → BP_PlayerState", False)],
            note="BeginPlay에서 CurrentHP를 MaxHP로 초기화해야 시작 시 HP 바가 가득 찹니다.")

    mk_step(prs, "실습 3: PlayerState HP + 데미지",
            2, "TakeDamage 함수 구현",
            [(0, "BP_PlayerState → My Blueprint → 함수 추가: TakeDamage", False),
             (0, "입력 파라미터: DamageAmount (Float)", False),
             (0, "로직: CurrentHP = Clamp(CurrentHP - DamageAmount, 0, MaxHP)", False),
             (0, "IsAlive 함수 추가: Return (CurrentHP > 0)", False),
             (0, "컴파일 & 저장", False)],
            tip="Clamp 노드를 사용해 HP가 0 미만 / MaxHP 초과가 되지 않도록 막으세요.",
            note="순수 함수(Pure Function)로 만들면 노드에 실행 핀이 없어 깔끔합니다.")

    mk_step(prs, "실습 3: PlayerState HP + 데미지",
            3, "IA_DebugDamage 생성 및 매핑",
            [(0, "IA_DebugDamage 생성 (Value Type: Digital)", False),
             (0, "IMC_Default → Mappings 추가: G 키 → IA_DebugDamage", False),
             (0, "저장", False)],
            note="실제 게임에서는 적이 데미지를 주지만, 오늘은 키보드로 테스트합니다.")

    mk_step(prs, "실습 3: PlayerState HP + 데미지",
            4, "Character에서 DebugDamage 바인딩",
            [(0, "Week3Character.h: UPROPERTY UInputAction* DebugDamageAction;", False),
             (0, "Week3Character.cpp: void DebugDamage(const FInputActionValue& Value);", False),
             (0, "SetupPlayerInputComponent: EIC->BindAction(DebugDamageAction, Started, this, &DebugDamage);", False),
             (0, "DebugDamage 구현:", False),
             (1, "APlayerState* PS = GetPlayerState();", False),
             (1, "Cast<ABP_PlayerState_C>(PS)->TakeDamage(10.f);", False)],
            tip="Cast 실패에 대비해 nullptr 체크를 추가하세요.",
            note="C++에서 BP 클래스를 캐스팅할 때는 헤더가 없으므로 BP 함수를 직접 호출하려면 BlueprintCallable 함수를 C++ 베이스 클래스에 선언해야 합니다.")

    mk_step(prs, "실습 3: PlayerState HP + 데미지",
            5, "빌드 · 테스트 · HP 바 연동 확인",
            [(0, "에디터 빌드 (Ctrl+Alt+F11 또는 빌드 버튼)", False),
             (0, "PIE 실행 → G 키 누르기", False),
             (0, "기대 결과: HP 바 Percent가 줄어든다", False),
             (0, "HP 0 도달 후 IsAlive 체크 → 향후 Death 로직 연결 포인트", False),
             (0, "Output Log: 에러 없음 확인", False)],
            tip="HP 바가 줄지 않으면: 바인딩 함수가 PlayerState를 올바르게 참조하는지 확인",
            note="이 시점에서 전체 실습 흐름 (Controller → Input → PlayerState → UI)이 연결됩니다.")

    # ── 실습 3 결과 확인 ─────────────────────────────────────────────────────
    mk_teal(prs, "실습 3 체크포인트 — 전체 흐름 완성",
            [(0, "✅ G 키 → IA_DebugDamage → Character → PlayerState.TakeDamage(10)", False),
             (0, "✅ HP 바가 10씩 줄어드는 것이 화면에 보인다", False),
             (0, "✅ HP가 0이 되면 IsAlive()가 false를 반환한다", False),
             (0, "✅ 전체 데이터 흐름: Input → Controller → Pawn → PlayerState → UI", False),
             (0, "▶ 다음 주: 이 기반 위에 적 스폰 + 피격 이벤트 추가", False)],
            note="흐름 전체를 칠판에 그리며 정리해 주세요. 학생들이 큰 그림을 보는 순간입니다.")

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION: 보너스
    # ═══════════════════════════════════════════════════════════════════════
    mk_section(prs, "보너스 — 탑다운 카메라 전환",
               "시간이 남으면 도전! (5분)",
               note="보너스는 선택 사항입니다. 실습이 밀리면 생략합니다.")

    mk_bullets(prs, "탑다운 카메라 — 빠른 설정",
               [(0, "BP_ThirdPersonCharacter → 뷰포트", False),
                (0, "SpringArm → TargetArmLength: 1200", False),
                (0, "SpringArm → Rotation: Pitch = -70, Yaw = 0, Roll = 0", False),
                (0, "SpringArm → Use Pawn Control Rotation: false 해제", False),
                (0, "Camera → Use Pawn Control Rotation: false", False),
                (0, "플레이 → 탑다운 시점 확인", False),
                (1, "VS-라이크 느낌! 마우스 조준은 다음 주 과제", False)],
               note="카메라 각도를 직접 조정해보는 재미있는 보너스입니다.")

    # ═══════════════════════════════════════════════════════════════════════
    # SECTION: 정리 & Q&A
    # ═══════════════════════════════════════════════════════════════════════
    mk_section(prs, "정리 & Q&A",
               "오늘 배운 것을 한 줄로 말할 수 있나요?",
               note="학생 한 명씩 오늘 배운 것 한 가지를 말하게 하면 복습 효과가 좋습니다.")

    # ── Slide 47: 오늘 배운 것 요약 ─────────────────────────────────────────
    mk_bullets(prs, "오늘 배운 것",
               [(0, "PlayerController — 입력·UI·빙의 담당, 폰과 독립적으로 존재", False),
                (0, "PlayerState — 플레이어 데이터 중앙 저장소, 폰 파괴 후에도 유지", False),
                (0, "Enhanced Input — IA + IMC 분리로 유연한 키 리맵핑", False),
                (0, "UMG ProgressBar + Percent Binding으로 실시간 HP 바", False),
                (0, "전체 흐름: Input → Controller → Pawn → PlayerState → UI", False)],
               note="핵심 키워드를 판서하면서 마무리하면 기억에 남습니다.")

    # ── Slide 48: 다음 주 예고 ──────────────────────────────────────────────
    mk_teal(prs, "다음 주 예고 (5주차 — 5/4)",
            [(0, "적(Enemy) 블루프린트 생성 및 랜덤 스폰", False),
             (0, "피격 이벤트 → PlayerState.TakeDamage 연결", False),
             (0, "UMG 심화: 웨이브 카운터, 킬 카운터 추가", False),
             (0, "게임오버 화면 구현", False),
             (0, "숙제: 현재 HP를 화면에 숫자로도 표시해 보기", False),
             (0, "질문·이슈는 스터디 단톡방에 올려주세요!", False)],
            note="숙제를 내주면 다음 주 시작 전에 학생들이 스스로 UMG를 더 탐구합니다.")

    prs.save(out_path)
    print(f"Saved: {out_path}")

if __name__ == "__main__":
    out = "/home/pjg_root/projects/Unreal/UnrealStudy/course_materials/언리얼_스터디_4주차_자료.pptx"
    build(out)
