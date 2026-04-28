"""
언리얼 스터디 4주차 자료 v2 생성 스크립트
원본 PPT를 복사한 뒤, 실습 0 (Possess 체험) 슬라이드 5장을 삽입하고
관련 슬라이드를 수정한다.
"""

import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── 경로 & 색상 상수 ─────────────────────────────────────
SRC = '/home/pjg_root/projects/Unreal/UnrealStudy/course_materials/언리얼_스터디_4주차_자료.pptx'
DST = '/home/pjg_root/projects/Unreal/UnrealStudy/course_materials/언리얼_스터디_4주차_자료_v2.pptx'

BLUE      = RGBColor(0x1B, 0x6B, 0x9A)
DARK_BLUE = RGBColor(0x0E, 0x45, 0x65)
BLACK     = RGBColor(0x00, 0x00, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x55, 0x55, 0x55)
FONT      = '맑은 고딕'

ANS = 'http://schemas.openxmlformats.org/drawingml/2006/main'


# ── XML 헬퍼 ─────────────────────────────────────────────

def _el(tag):
    return etree.Element(f'{{{ANS}}}{tag}')


def _sub(parent, tag):
    return etree.SubElement(parent, f'{{{ANS}}}{tag}')


def make_para_el(text, size_hundredths, bold=False, color_hex='000000',
                 align=None):
    """단락 <a:p> 엘리먼트를 생성한다. 텍스트는 lxml이 자동 이스케이프."""
    p = _el('p')

    # 단락 정렬
    if align is not None:
        pPr = _sub(p, 'pPr')
        align_map = {PP_ALIGN.CENTER: 'ctr', PP_ALIGN.LEFT: 'l',
                     PP_ALIGN.RIGHT: 'r'}
        pPr.set('algn', align_map.get(align, 'l'))

    r = _sub(p, 'r')
    rPr = _sub(r, 'rPr')
    rPr.set('lang', 'ko-KR')
    rPr.set('sz', str(size_hundredths))
    rPr.set('b', '1' if bold else '0')
    rPr.set('dirty', '0')

    solidFill = _sub(rPr, 'solidFill')
    srgbClr = _sub(solidFill, 'srgbClr')
    srgbClr.set('val', color_hex)

    latin = _sub(rPr, 'latin')
    latin.set('typeface', FONT)

    t = _sub(r, 't')
    t.text = text

    return p


def get_txBody(tf):
    return tf._txBody


def clear_txBody(tf):
    txBody = get_txBody(tf)
    for p in txBody.findall(qn('a:p')):
        txBody.remove(p)


def append_para(tf, text, size_pt, bold=False, color=None, align=None):
    color_hex = '{:06X}'.format(
        (color.rgb if hasattr(color, 'rgb') else 0) if color else 0
    )
    if color is None:
        color_hex = '000000'
    elif isinstance(color, RGBColor):
        color_hex = '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])

    el = make_para_el(
        text,
        int(size_pt * 100),
        bold=bold,
        color_hex=color_hex,
        align=align,
    )
    get_txBody(tf).append(el)


def insert_para_before_first(tf, text, size_pt, bold=False, color=None):
    color_hex = '000000'
    if color is not None and isinstance(color, RGBColor):
        color_hex = '{:02X}{:02X}{:02X}'.format(color[0], color[1], color[2])
    el = make_para_el(text, int(size_pt * 100), bold=bold, color_hex=color_hex)
    txBody = get_txBody(tf)
    existing = txBody.findall(qn('a:p'))
    if existing:
        txBody.insert(list(txBody).index(existing[0]), el)
    else:
        txBody.append(el)


# ── 도형 생성 헬퍼 ────────────────────────────────────────

def add_textbox(slide, left_in, top_in, w_in, h_in):
    return slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(w_in), Inches(h_in)
    )


def add_rect(slide, left_in, top_in, w_in, h_in, fill_color: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left_in), Inches(top_in),
        Inches(w_in), Inches(h_in)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def get_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name == 'Blank':
            return layout
    return prs.slide_layouts[6]


# ── 슬라이드 레이아웃 빌더 ───────────────────────────────

def build_section_title(slide, title, subtitle):
    """섹션 표지: 가운데 큰 제목 + 회색 소제목"""
    tb1 = add_textbox(slide, 0.50, 2.40, 12.33, 2.00)
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    clear_txBody(tf1)
    append_para(tf1, title, 44, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    tb2 = add_textbox(slide, 0.50, 4.50, 12.33, 0.80)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    clear_txBody(tf2)
    append_para(tf2, subtitle, 22, color=GRAY, align=PP_ALIGN.CENTER)


def build_step_slide(slide, section_label, step_num, step_title,
                     bullets, warning=None):
    """실습 Step 슬라이드"""
    # 섹션 레이블
    tb1 = add_textbox(slide, 0.50, 0.15, 12.30, 0.70)
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    clear_txBody(tf1)
    append_para(tf1, section_label, 24, bold=True, color=GRAY, align=PP_ALIGN.LEFT)

    # Step 배지
    add_rect(slide, 0.50, 0.95, 1.50, 0.65, BLUE)
    tb3 = add_textbox(slide, 0.50, 0.95, 1.50, 0.65)
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    clear_txBody(tf3)
    append_para(tf3, f'Step {step_num}', 19, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)

    # Step 제목
    tb4 = add_textbox(slide, 2.10, 0.95, 10.70, 0.65)
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    clear_txBody(tf4)
    append_para(tf4, step_title, 22, bold=True, color=BLACK, align=PP_ALIGN.LEFT)

    # 구분선
    add_rect(slide, 0.50, 1.65, 12.30, 0.04, BLUE)

    # 본문
    tb6 = add_textbox(slide, 0.50, 1.80, 12.30, 4.60)
    tf6 = tb6.text_frame
    tf6.word_wrap = True
    clear_txBody(tf6)
    for text, size_pt in bullets:
        append_para(tf6, text, size_pt, color=BLACK)

    # 하단 경고
    if warning:
        add_rect(slide, 0.50, 6.60, 12.30, 0.70, DARK_BLUE)
        tb8 = add_textbox(slide, 0.65, 6.60, 12.00, 0.70)
        tf8 = tb8.text_frame
        tf8.word_wrap = True
        clear_txBody(tf8)
        append_para(tf8, warning, 16, color=WHITE, align=PP_ALIGN.LEFT)


def build_callout(slide, title, bullets):
    """파란 배경 콜아웃 (체크포인트/AHA/예고)
    bullets: list of (text, size_pt, bold)
    """
    add_rect(slide, 0.15, 0.15, 13.03, 7.20, BLUE)
    add_rect(slide, 0.35, 0.35, 12.63, 6.80, DARK_BLUE)

    tb3 = add_textbox(slide, 0.55, 0.40, 12.20, 1.00)
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    clear_txBody(tf3)
    append_para(tf3, title, 30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    tb4 = add_textbox(slide, 0.90, 1.50, 11.80, 5.50)
    tf4 = tb4.text_frame
    tf4.word_wrap = True
    clear_txBody(tf4)
    for text, size_pt, bold in bullets:
        append_para(tf4, text, size_pt, bold=bold, color=WHITE)


# ── 기존 슬라이드 수정 ────────────────────────────────────

def modify_slide_goal(slides):
    """슬라이드 2: 강의 목표 맨 앞에 실습 0 항목 추가"""
    slide = slides[1]
    for shape in slide.shapes:
        if shape.name == 'TextBox 3' and shape.has_text_frame:
            insert_para_before_first(
                shape.text_frame,
                '• 실습 0: C++ PlayerController로 원하는 폰에 직접 빙의할 수 있다',
                21, color=BLACK
            )
    print('  [OK] slide 2: 목표 수정')


def modify_slide_timetable(slides):
    """슬라이드 4: 시간표 90분 기준 재배분 + 실습 0 추가"""
    slide = slides[3]
    for shape in slide.shapes:
        if shape.name == 'TextBox 1' and shape.has_text_frame:
            tf = shape.text_frame
            clear_txBody(tf)
            append_para(tf, '오늘 시간 배분 (90분)', 34,
                        bold=True, color=BLACK, align=PP_ALIGN.LEFT)

        if shape.name == 'TextBox 3' and shape.has_text_frame:
            tf = shape.text_frame
            clear_txBody(tf)
            rows = [
                ('• 에디터 기초 숙달 (15분)', 21),
                ('    • 뷰포트 조작 / 기즈모 / 플레이 모드', 18),
                ('• 이론 A — PlayerController (15분)', 21),
                ('• 실습 0 — Possess 체험 (10분)  [신규]', 21),
                ('• 이론 B — PlayerState (8분)', 21),
                ('• 이론 C — Enhanced Input 심화 (10분)', 21),
                ('• 실습 1: 스프린트 구현 (10분)', 21),
                ('• 실습 2: HP 바 위젯 (10분)', 21),
                ('• 실습 3: PlayerState HP + 데미지 (7분)', 21),
                ('• 보너스 / Q&A (5분+)', 21),
            ]
            for text, size_pt in rows:
                append_para(tf, text, size_pt, color=BLACK)
    print('  [OK] slide 4: 시간표 수정')


def modify_slide_demo(slides):
    """슬라이드 16: 데모 마지막에 실습 0 브릿지 추가"""
    slide = slides[15]
    for shape in slide.shapes:
        if shape.name == 'TextBox 3' and shape.has_text_frame:
            tf = shape.text_frame
            append_para(tf, '', 18, color=BLACK)
            append_para(
                tf,
                '▶ 이제 직접 체험해보자 — 실습 0',
                21, bold=True, color=BLUE
            )
    print('  [OK] slide 16: 브릿지 추가')


def modify_slide_lab2_step3(slides):
    """슬라이드 35 (삽입 후 p.40): 실습 2 Step 3에 실습 0 확장 명시"""
    slide = slides[34]
    for shape in slide.shapes:
        if shape.name == 'TextBox 6' and shape.has_text_frame:
            insert_para_before_first(
                shape.text_frame,
                '    ※ 실습 0에서 만든 Week3PlayerController를 확장합니다',
                17, bold=True, color=BLUE
            )
    print('  [OK] slide 35: 실습 2 Step 3 수정')


def modify_slide_summary(slides):
    """슬라이드 47 (삽입 후 p.52): 실습 0 bullet 추가 + PlayerController 강조"""
    slide = slides[46]
    for shape in slide.shapes:
        if shape.name == 'TextBox 3' and shape.has_text_frame:
            tf = shape.text_frame

            # 첫 번째 단락(PlayerController) 텍스트 교체 + 색상을 파란색으로
            txBody = get_txBody(tf)
            first_p = txBody.findall(qn('a:p'))[0]
            for r in first_p.findall(qn('a:r')):
                t = r.find(qn('a:t'))
                if t is not None:
                    t.text = ('▶ [기반] PlayerController — Possess()로 빙의, '
                              '입력·UI·빙의 담당 '
                              '(실습 0→1→2→3의 핵심)')
                rPr = r.find(qn('a:rPr'))
                if rPr is not None:
                    sf = rPr.find(qn('a:solidFill'))
                    if sf is None:
                        sf = _sub(rPr, 'solidFill')
                    clr = sf.find(qn('a:srgbClr'))
                    if clr is None:
                        clr = _sub(sf, 'srgbClr')
                    clr.set('val', '1B6B9A')
                    rPr.set('b', '1')

            # 실습 0 bullet 맨 앞에 삽입
            insert_para_before_first(
                tf,
                '• 실습 0: PlayerController::Possess()로 원하는 폰에 직접 빙의',
                21, color=BLACK
            )
    print('  [OK] slide 47: 달성 요약 수정')


# ── 신규 슬라이드 생성 ────────────────────────────────────

def add_slides(prs):
    layout = get_blank_layout(prs)

    # p.17 섹션 표지
    s17 = prs.slides.add_slide(layout)
    build_section_title(s17,
                        '실습 0 — Possess: 두 개의 몸',
                        '두 개의 몸, 어느 쪽에 들어갈지 우리가 결정한다')
    print('  [OK] 신규 p.17: 섹션 표지')

    # p.18 Step 1: 맵 준비
    s18 = prs.slides.add_slide(layout)
    build_step_slide(
        s18,
        section_label='실습 0: Possess 체험',
        step_num=1,
        step_title='맵 준비 — 캐릭터 A·B 배치',
        bullets=[
            ('• Content Browser → BP_ThirdPersonCharacter 두 개 드래그 → 레벨에 배치', 20),
            ('• World Outliner에서 이름 변경: CharacterA, CharacterB', 20),
            ('• World Settings → Game Mode Override → BP_Gamemode 선택', 20),
            ('• BP_Gamemode 열기 → Default Pawn Class: None 으로 변경', 20),
            ('• 저장 (Ctrl+S)', 20),
        ],
        warning='★ DefaultPawnClass=None 이어야 게임 시작 시 자동으로 폰이 스폰되지 않습니다.'
    )
    print('  [OK] 신규 p.18: Step 1')

    # p.19 Step 2: PossessSelector 생성
    s19 = prs.slides.add_slide(layout)
    build_step_slide(
        s19,
        section_label='실습 0: Possess 체험',
        step_num=2,
        step_title='AWeek3PossessSelector C++ 클래스 생성',
        bullets=[
            ('• Tools → New C++ Class → Actor → 이름: Week3PossessSelector', 20),
            ('• Week3PossessSelector.h 에 변수 추가:', 20),
            ('    • UPROPERTY(EditAnywhere, Category="Possess") APawn* CharacterToControl;', 17),
            ('• 에디터 빌드 (Ctrl+Alt+F11)', 20),
            ('• 레벨에 BP_Week3PossessSelector 배치', 20),
            ('• 디테일 패널 → CharacterToControl → CharacterA 선택', 20),
        ],
        warning='★ EditAnywhere 이므로 레벨의 디테일 패널에서 폰을 자유롭게 교체할 수 있습니다.'
    )
    print('  [OK] 신규 p.19: Step 2')

    # p.20 Step 3: PlayerController 생성
    s20 = prs.slides.add_slide(layout)
    build_step_slide(
        s20,
        section_label='실습 0: Possess 체험',
        step_num=3,
        step_title='AWeek3PlayerController 생성 + BeginPlay에서 Possess',
        bullets=[
            ('• Tools → New C++ Class → PlayerController → 이름: Week3PlayerController', 20),
            ('• BeginPlay 오버라이드: TActorIterator<AWeek3PossessSelector>로 Selector 탐색', 20),
            ('• Selector->CharacterToControl 가져와서 Possess(target) 호출', 20),
            ('• BP_Gamemode → Player Controller Class: BP_Week3PlayerController 설정', 20),
            ('• 빌드 (Ctrl+Alt+F11) 후 PIE 실행', 20),
        ],
        warning='★ TActorIterator 결과가 없으면 Selector가 레벨에 제대로 배치됐는지 확인하세요.'
    )
    print('  [OK] 신규 p.20: Step 3')

    # p.21 체크포인트 + AHA
    s21 = prs.slides.add_slide(layout)
    build_callout(
        s21,
        title='실습 0 체크포인트 + AHA!',
        bullets=[
            ('• ✅ PIE 시작 시 CharacterToControl 에 지정한 폰으로 빙의된다', 20, False),
            ('• ✅ 디테일 패널에서 CharacterB 로 바꾸면 다음 PIE에 B로 빙의된다', 20, False),
            ('• ✅ Week3Character 코드를 단 한 줄도 수정하지 않았다', 20, False),
            ('', 14, False),
            ('• 💡 AHA! "코드 한 줄 안 바꾸고 변수만 바꿨는데 다른 몸으로 들어갔다"', 20, True),
            ('', 14, False),
            ('• ▶ 이 PlayerController를 실습 1(스프린트) · 2(HP 바) · 3(데미지)의 기반으로 확장한다', 20, False),
        ]
    )
    print('  [OK] 신규 p.21: 체크포인트 + AHA')


def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    els = list(xml_slides)
    el = els[old_index]
    xml_slides.remove(el)
    xml_slides.insert(new_index, el)


# ── 메인 ─────────────────────────────────────────────────

def main():
    print('원본 복사...')
    shutil.copy2(SRC, DST)
    prs = Presentation(DST)
    slides = list(prs.slides)
    print(f'  원본 슬라이드 수: {len(slides)}')

    print('\n[1] 기존 슬라이드 수정')
    modify_slide_goal(slides)
    modify_slide_timetable(slides)
    modify_slide_demo(slides)
    modify_slide_lab2_step3(slides)
    modify_slide_summary(slides)

    print('\n[2] 신규 슬라이드 5장 추가')
    add_slides(prs)
    total = len(list(prs.slides))  # 53
    print(f'  총 슬라이드 수: {total}')

    print('\n[3] 신규 슬라이드를 p.17~21 위치로 이동')
    insert_at = 16  # 0-indexed (p.17 = index 16)
    for i in range(5):
        old_idx = total - 5 + i
        new_idx = insert_at + i
        move_slide(prs, old_idx, new_idx)
        print(f'  idx {old_idx} -> {new_idx}')

    print(f'\n저장: {DST}')
    prs.save(DST)
    print('완료!')
    print(f'최종 슬라이드 수: {len(list(prs.slides))}')


if __name__ == '__main__':
    main()
